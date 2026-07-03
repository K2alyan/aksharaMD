from __future__ import annotations

import logging
from collections import Counter
from pathlib import Path

logger = logging.getLogger(__name__)

from ...context import CompilationContext
from ...models.asset import Asset
from ...models.block import Block, BlockType
from ...models.document import Document
from ..base import ParserPlugin
from ..registry import register_parser

_PT_TO_PX = 4 / 3  # 1 pt = 1.333 px (irrelevant here — we just compare pt values)


def _para_font_size_pt(para) -> float | None:
    """Return the dominant font size of a paragraph in points, or None."""
    try:
        sizes = []
        for run in para.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
        if sizes:
            return max(sizes)
        # Fall back to paragraph-level font
        if para.font.size:
            return para.font.size.pt
    except Exception:
        pass
    return None


def _para_has_explicit_bullet(para) -> bool:
    """Return True if the paragraph has explicit bullet formatting in its own XML."""
    try:
        from pptx.oxml.ns import qn
        pPr = para._p.find(qn("a:pPr"))
        if pPr is None:
            return False
        if pPr.find(qn("a:buNone")) is not None:
            return False
        if (pPr.find(qn("a:buChar")) is not None or
                pPr.find(qn("a:buAutoNum")) is not None or
                pPr.find(qn("a:buFont")) is not None):
            return True
    except Exception:
        pass
    return False


def _para_level(para) -> int:
    """Return the indentation level of a paragraph (0 = top level)."""
    try:
        from pptx.oxml.ns import qn
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            lvl = pPr.get("lvl")
            if lvl is not None:
                return max(0, int(lvl))
    except Exception:
        pass
    return 0


def _is_body_placeholder(shape) -> bool:
    """Return True if this shape is a body/content placeholder (layout-inherited bullets)."""
    try:
        from pptx.enum.placeholders import PP_PLACEHOLDER
        ph = shape.placeholder_format
        if ph is not None and ph.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return True
    except Exception:
        pass
    return False


def _text_frame_is_list(paras: list, shape=None) -> bool:
    """
    Return True if this text frame is a bullet list.

    Checks in order:
    1. Explicit bullet XML markers
    2. Any paragraph with indent level > 0 (layout-inherited nested bullet)
    3. Body/content placeholder type (layout-inherited top-level bullets)
    4. Content heuristic: multiple short, non-prose paragraphs
    """
    if not paras:
        return False

    try:
        from pptx.oxml.ns import qn
        for p in paras:
            pPr = p._p.find(qn("a:pPr"))
            if pPr is not None:
                if pPr.find(qn("a:buNone")) is None and (
                    pPr.find(qn("a:buChar")) is not None or
                    pPr.find(qn("a:buAutoNum")) is not None
                ):
                    return True
                lvl = pPr.get("lvl")
                if lvl is not None and int(lvl) > 0:
                    return True
    except Exception:
        pass

    # Body/content placeholder = layout-inherited bullets by convention
    if shape is not None and _is_body_placeholder(shape):
        return True

    # Heuristic: 2+ short, non-prose paragraphs
    if len(paras) < 2:
        return False
    for p in paras:
        text = p.text.strip()
        words = text.split()
        if len(words) > 30:
            return False
        if text.count(". ") >= 3:
            return False
    return True


def _table_to_markdown(table) -> str:
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(rows)


def _shape_to_blocks(shape, slide_num: int, idx: int,
                     title_shape) -> tuple[list[Block], int]:
    """
    Convert a single PPTX shape into one or more Blocks.

    Strategy:
    - Tables → TABLE block
    - Text frames → analyse each paragraph:
        * Body/content placeholder or layout-inherited bullets → LIST block with indentation
        * Large font (>= 20 pt) and few words → HEADING level 3
        * Otherwise → PARAGRAPH block per logical group
    """
    blocks: list[Block] = []

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            md = _table_to_markdown(shape.table)
            if md:
                blocks.append(Block(type=BlockType.TABLE, content=md,
                                    page=slide_num, index=idx))
            return blocks, idx + len(blocks)

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return blocks, idx  # no text to extract

        if not shape.has_text_frame:
            return blocks, idx

    except Exception:
        return blocks, idx

    # --- text frame ---
    tf = shape.text_frame
    paras = [p for p in tf.paragraphs if p.text.strip()]
    if not paras:
        return blocks, idx

    # Decide upfront if the whole text frame is a list (pass shape for placeholder check)
    frame_is_list = _text_frame_is_list(paras, shape=shape)

    if frame_is_list:
        items = []
        for p in paras:
            text = p.text.strip()
            if text:
                lvl = _para_level(p)
                indent = "  " * lvl
                items.append(f"{indent}- {text}")
        content = "\n".join(items)
        blocks.append(Block(type=BlockType.LIST, content=content,
                            page=slide_num, index=idx))
        idx += 1
    else:
        # Mixed content — process paragraph by paragraph
        pending_flush: list[str] = []

        def flush_pending():
            nonlocal idx
            if pending_flush:
                text = "\n".join(pending_flush)
                if text:
                    blocks.append(Block(type=BlockType.PARAGRAPH, content=text,
                                        page=slide_num, index=idx))
                    idx += 1
                pending_flush.clear()

        for para in paras:
            text = para.text.strip()
            if not text:
                continue

            size = _para_font_size_pt(para)
            is_big = size is not None and size >= 20
            is_short = len(text.split()) <= 12
            has_explicit_bullet = _para_has_explicit_bullet(para)

            if has_explicit_bullet:
                flush_pending()
                lvl = _para_level(para)
                indent = "  " * lvl
                blocks.append(Block(type=BlockType.LIST, content=f"{indent}- {text}",
                                    page=slide_num, index=idx))
                idx += 1
            elif is_big and is_short:
                flush_pending()
                blocks.append(Block(type=BlockType.HEADING, content=text,
                                    level=3, page=slide_num, index=idx))
                idx += 1
            else:
                pending_flush.append(text)

        flush_pending()

    return blocks, idx


class PptxParser(ParserPlugin):
    name = "pptx_parser"
    supported_types = ["pptx"]

    def execute(self, ctx: CompilationContext) -> CompilationContext:
        from pptx import Presentation

        path = Path(ctx.source)
        try:
            prs = Presentation(str(path))
        except Exception as e:
            ctx.error("PPTX_PARSE_ERROR", str(e))
            return ctx

        blocks: list[Block] = []
        assets: list[Asset] = []
        idx = 0
        title: str | None = None
        image_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_title_text = None
            title_shape = slide.shapes.title

            # Slide heading
            if title_shape:
                slide_title_text = title_shape.text_frame.text.strip() if title_shape.has_text_frame else ""
                if slide_title_text:
                    if not title:
                        title = slide_title_text
                    blocks.append(Block(
                        type=BlockType.HEADING, content=slide_title_text,
                        level=2, page=slide_num, index=idx,
                    ))
                    idx += 1

            # Count images, process all other shapes
            for shape in slide.shapes:
                if shape == title_shape:
                    continue
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                        try:
                            img_bytes = shape.image.blob
                            if img_bytes:
                                asset_id = f"img_{slide_num}_{idx}"
                                assets.append(Asset(
                                    id=asset_id, type="image",
                                    image_bytes=img_bytes,
                                    metadata={"slide": slide_num},
                                ))
                                alt = shape.name or ""
                                blocks.append(Block(
                                    type=BlockType.IMAGE, content=alt,
                                    page=slide_num, index=idx,
                                    metadata={"asset_id": asset_id},
                                ))
                                idx += 1
                        except Exception:
                            pass
                        continue
                except Exception:
                    pass

                new_blocks, idx = _shape_to_blocks(shape, slide_num, idx, title_shape)
                blocks.extend(new_blocks)

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(Block(
                        type=BlockType.BLOCKQUOTE,
                        content=f"[Speaker notes] {notes}",
                        page=slide_num, index=idx,
                    ))
                    idx += 1

        # Deduplicate bullet items that repeat verbatim across 3+ slides.
        # Template decks often carry the same boilerplate bullets on every slide;
        # keeping only the first occurrence preserves meaning without repetition.
        bullet_counts: Counter[str] = Counter()
        for block in blocks:
            if block.type == BlockType.LIST:
                for line in block.content.splitlines():
                    key = line.strip()
                    if key.startswith("- "):
                        bullet_counts[key] += 1

        if any(v >= 3 for v in bullet_counts.values()):
            seen_bullets: set[str] = set()
            deduped: list[Block] = []
            for block in blocks:
                if block.type == BlockType.LIST:
                    kept = []
                    for line in block.content.splitlines():
                        key = line.strip()
                        if key.startswith("- ") and bullet_counts[key] >= 3:
                            if key not in seen_bullets:
                                seen_bullets.add(key)
                                kept.append(line)
                        else:
                            kept.append(line)
                    if kept:
                        deduped.append(block.model_copy(update={"content": "\n".join(kept)}))
                else:
                    deduped.append(block)
            blocks = deduped

        ctx.document = Document(
            source=str(path),
            file_type="pptx",
            title=title,
            pages=len(prs.slides),
            blocks=blocks,
            assets=assets,
            metadata={"slides": len(prs.slides), "images": image_count},
        ).compute_id()
        return ctx


register_parser("pptx", PptxParser)
