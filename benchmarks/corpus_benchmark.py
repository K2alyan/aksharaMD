#!/usr/bin/env python3
"""
Multi-tool benchmark runner against the full AksharaMD benchmark corpus.

Tools compared:
  - AksharaMD       (this project)
  - MarkItDown     (Microsoft)
  - Docling        (IBM)
  - Unstructured   (Unstructured.io)
  - PyMuPDF4LLM   (pymupdf4llm, PDF only)
  - pypdf          (PDF only, text extraction)

Usage:
  python -m benchmarks.corpus_benchmark --corpus-dir ../Downloads/benchmark_corpus
  python -m benchmarks.corpus_benchmark --corpus-dir ../Downloads/benchmark_corpus --tools omnimark markitdown
  python -m benchmarks.corpus_benchmark --corpus-dir ../Downloads/benchmark_corpus --types pdf docx --limit 10
"""
from __future__ import annotations

import argparse
import csv
import importlib
import json

# Suppress noisy logging from docling / RapidOCR / transformers
import logging
import os
import sys
import tempfile
import time
from dataclasses import asdict, dataclass
from pathlib import Path

from benchmarks.metrics import compute_metrics

os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")
for noisy in ("docling", "docling_core", "rapidocr", "transformers",
              "huggingface_hub", "accelerate", "torch"):
    logging.getLogger(noisy).setLevel(logging.ERROR)

# ── tool support matrix ───────────────────────────────────────────────────────
# Maps file extension → set of tool names that can handle it.
TOOL_SUPPORT: dict[str, set[str]] = {
    ".pdf":   {"aksharamd", "naive", "markitdown", "docling", "unstructured", "pymupdf4llm", "pypdf", "pdfplumber"},
    ".docx":  {"aksharamd", "naive", "markitdown", "docling", "unstructured"},
    ".doc":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".pptx":  {"aksharamd", "naive", "markitdown", "docling", "unstructured"},
    ".ppt":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".xlsx":  {"aksharamd", "naive", "markitdown", "unstructured"},
    ".xls":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".msg":   {"aksharamd", "naive", "markitdown"},
    ".eml":   {"aksharamd", "naive", "markitdown"},
    ".epub":  {"aksharamd", "naive", "markitdown", "unstructured"},
    ".odt":   {"aksharamd"},
    ".ods":   {"aksharamd"},
    ".odp":   {"aksharamd"},
    ".rtf":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".txt":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".md":    {"aksharamd", "naive", "markitdown"},
    ".html":  {"aksharamd", "naive", "markitdown", "docling", "unstructured"},
    ".htm":   {"aksharamd", "naive", "markitdown", "docling", "unstructured"},
    ".csv":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".json":  {"aksharamd", "naive", "markitdown"},
    ".jsonl": {"aksharamd", "naive", "markitdown"},
    ".xml":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".rss":   {"aksharamd", "naive", "markitdown"},
    ".atom":  {"aksharamd", "naive", "markitdown"},
    ".zip":   {"aksharamd", "naive", "markitdown", "unstructured"},
    ".tar":   {"aksharamd"},
    ".gz":    {"aksharamd"},
    ".7z":    {"aksharamd"},
    ".ipynb": {"aksharamd", "naive", "markitdown", "unstructured"},
    ".jpg":   {"aksharamd", "markitdown", "docling", "unstructured"},
    ".jpeg":  {"aksharamd", "markitdown", "docling", "unstructured"},
    ".png":   {"aksharamd", "markitdown", "docling", "unstructured"},
    ".mp3":   {"aksharamd", "markitdown"},
    ".wav":   {"aksharamd", "markitdown"},
    ".m4a":   {"aksharamd", "markitdown"},
    ".mp4":   {"aksharamd", "markitdown"},
}

ALL_TOOLS = ["naive", "aksharamd", "markitdown", "docling", "unstructured", "pymupdf4llm", "pypdf", "pdfplumber"]

FILE_TYPE_MAP = {
    "pdf": [".pdf"],
    "pptx": [".pptx"], "ppt": [".ppt"],
    "docx": [".docx"], "doc": [".doc"],
    "xlsx": [".xlsx"], "xls": [".xls"],
    "msg": [".msg"],
    "epub": [".epub"],
    "txt": [".txt"], "md": [".md", ".markdown"],
    "html": [".html", ".htm"],
    "csv": [".csv"],
    "json": [".json"], "jsonl": [".jsonl"],
    "xml": [".xml"], "rss": [".rss"], "atom": [".atom"],
    "zip": [".zip"],
    "ipynb": [".ipynb"],
    "jpg": [".jpg", ".jpeg"], "png": [".png"],
    "mp3": [".mp3"], "wav": [".wav"], "m4a": [".m4a"], "mp4": [".mp4"],
    "youtube_url": [".url.txt"], "wikipedia_url": [],
}

TIMEOUT_SECONDS = 120


# ── per-tool runners ──────────────────────────────────────────────────────────

def _run_aksharamd(path: Path) -> tuple[str, float]:
    from aksharamd.compiler import Compiler
    with tempfile.TemporaryDirectory() as tmp:
        compiler = Compiler(output_dir=tmp)
        t0 = time.perf_counter()
        ctx = compiler.compile(str(path))
        elapsed = time.perf_counter() - t0
        md_path = Path(tmp) / "document.md"
        text = md_path.read_text(encoding="utf-8") if md_path.exists() else ""
        if not text and ctx.document:
            text = "\n\n".join(b.content for b in ctx.document.blocks)
        return text, elapsed


def _run_markitdown(path: Path) -> tuple[str, float]:
    from markitdown import MarkItDown
    md = MarkItDown()
    t0 = time.perf_counter()
    result = md.convert(str(path))
    elapsed = time.perf_counter() - t0
    return result.text_content or "", elapsed


def _run_docling(path: Path) -> tuple[str, float]:
    from docling.document_converter import DocumentConverter
    converter = DocumentConverter()
    t0 = time.perf_counter()
    result = converter.convert(str(path))
    elapsed = time.perf_counter() - t0
    text = result.document.export_to_markdown()
    return text or "", elapsed


def _unstructured_table_to_md(html: str) -> str:
    """Convert unstructured's text_as_html table to Markdown."""
    try:
        from bs4 import BeautifulSoup
        rows = BeautifulSoup(html, "html.parser").find_all("tr")
        if not rows:
            return ""
        md_rows: list[str] = []
        for i, row in enumerate(rows):
            cells = [td.get_text(separator=" ", strip=True) for td in row.find_all(["td", "th"])]
            if not cells:
                continue
            md_rows.append("| " + " | ".join(c.replace("|", "\\|") for c in cells) + " |")
            if i == 0:
                md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        return "\n".join(md_rows)
    except Exception:
        return ""


def _run_unstructured(path: Path) -> tuple[str, float]:
    from unstructured.partition.auto import partition
    t0 = time.perf_counter()
    elements = partition(filename=str(path))
    elapsed = time.perf_counter() - t0

    parts: list[str] = []
    for el in elements:
        el_type = type(el).__name__
        text = str(el).strip()
        if not text:
            continue
        if el_type == "Title":
            parts.append(f"## {text}")
        elif el_type in ("Header", "SectionHeader"):
            parts.append(f"### {text}")
        elif el_type == "ListItem":
            parts.append(f"- {text}")
        elif el_type == "Table":
            html = getattr(getattr(el, "metadata", None), "text_as_html", None)
            md = _unstructured_table_to_md(html) if html else ""
            parts.append(md if md else text)
        elif el_type == "FigureCaption":
            parts.append(f"*{text}*")
        elif el_type in ("CodeSnippet", "Code"):
            parts.append(f"```\n{text}\n```")
        elif el_type in ("Footer", "PageBreak"):
            continue  # skip structural noise
        else:
            parts.append(text)  # NarrativeText, Address, Formula, etc.

    return "\n\n".join(parts), elapsed


def _run_pymupdf4llm(path: Path) -> tuple[str, float]:
    import pymupdf4llm
    t0 = time.perf_counter()
    text = pymupdf4llm.to_markdown(str(path))
    elapsed = time.perf_counter() - t0
    return text or "", elapsed


def _run_pypdf(path: Path) -> tuple[str, float]:
    from pypdf import PdfReader
    t0 = time.perf_counter()
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    elapsed = time.perf_counter() - t0
    return "\n\n".join(pages), elapsed


def _run_pdfplumber(path: Path) -> tuple[str, float]:
    import pdfplumber
    t0 = time.perf_counter()
    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            page_parts: list[str] = []
            text = page.extract_text()
            if text:
                page_parts.append(text)
            for table in (page.extract_tables() or []):
                if not table:
                    continue
                col_count = max(len(row) for row in table)
                rows_md = []
                for row in table:
                    cells = [str(c or "").replace("\n", " ").strip() for c in row]
                    cells += [""] * (col_count - len(cells))
                    rows_md.append("| " + " | ".join(cells) + " |")
                if rows_md:
                    sep = "| " + " | ".join(["---"] * col_count) + " |"
                    rows_md.insert(1, sep)
                    page_parts.append("\n".join(rows_md))
            if page_parts:
                parts.append("\n\n".join(page_parts))
    elapsed = time.perf_counter() - t0
    return "\n\n".join(parts), elapsed


def _run_naive(path: Path) -> tuple[str, float]:
    """Raw text extraction with zero cleanup — format-appropriate but no post-processing."""
    import time
    ext = path.suffix.lower()
    t0 = time.perf_counter()

    if ext == ".pdf":
        import fitz
        doc = fitz.open(str(path))
        text = "".join(page.get_text() for page in doc)
        doc.close()

    elif ext in (".txt", ".csv", ".json", ".jsonl", ".xml", ".rss", ".atom",
                 ".html", ".htm", ".md", ".rst", ".tex", ".yaml", ".toml", ".sh",
                 ".py", ".js", ".ts", ".go", ".rs", ".java", ".c", ".cpp", ".sql"):
        with open(path, encoding="utf-8", errors="replace") as f:
            text = f.read()

    elif ext == ".docx":
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        text = "\n".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, "text")
        )

    elif ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        parts.append(str(cell.value))
        text = "\n".join(parts)

    elif ext == ".epub":
        import ebooklib
        from bs4 import BeautifulSoup
        from ebooklib import epub
        book = epub.read_epub(str(path), options={"ignore_ncx": True})
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            parts.append(soup.get_text())
        text = "\n".join(parts)

    elif ext in (".eml", ".msg"):
        with open(path, encoding="utf-8", errors="replace") as f:
            text = f.read()

    elif ext == ".ipynb":
        with open(path, encoding="utf-8", errors="replace") as f:
            text = f.read()

    elif ext == ".zip":
        import zipfile
        try:
            with zipfile.ZipFile(str(path)) as zf:
                text = "\n".join(zf.namelist())
        except Exception:
            text = ""

    elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".webp", ".gif"):
        # No text layer — naive extraction not applicable
        text = ""

    elif ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac", ".mp4", ".webm"):
        text = ""

    elif ext in (".tar", ".gz", ".bz2", ".xz", ".7z", ".tgz"):
        text = ""

    else:
        with open(path, encoding="utf-8", errors="replace") as f:
            text = f.read()

    elapsed = time.perf_counter() - t0
    return text, elapsed


RUNNERS: dict[str, callable] = {
    "naive":        _run_naive,
    "aksharamd":    _run_aksharamd,
    "markitdown":   _run_markitdown,
    "docling":      _run_docling,
    "unstructured": _run_unstructured,
    "pymupdf4llm":  _run_pymupdf4llm,
    "pypdf":        _run_pypdf,
    "pdfplumber":   _run_pdfplumber,
}


# ── availability check ────────────────────────────────────────────────────────

def check_tool_available(tool: str) -> bool:
    MODULE_MAP = {
        "naive":        "fitz",
        "aksharamd":    "aksharamd.compiler",
        "markitdown":   "markitdown",
        "docling":      "docling.document_converter",
        "unstructured": "unstructured.partition.auto",
        "pymupdf4llm":  "pymupdf4llm",
        "pypdf":        "pypdf",
        "pdfplumber":   "pdfplumber",
    }
    try:
        importlib.import_module(MODULE_MAP[tool])
        return True
    except ImportError:
        return False


# ── result dataclass ──────────────────────────────────────────────────────────

@dataclass
class ToolResult:
    file: str
    file_type: str
    file_size_kb: float
    tool: str
    supported: bool
    available: bool
    status: str            # "ok" | "unsupported" | "unavailable" | "error" | "timeout"
    error: str = ""
    token_count: int = 0
    char_count: int = 0
    heading_count: int = 0
    table_count: int = 0
    noise_line_count: int = 0
    duplicate_line_count: int = 0
    avg_paragraph_tokens: float = 0.0
    quality_score: int = 0
    elapsed_seconds: float = 0.0


# ── main runner ───────────────────────────────────────────────────────────────

def run_file(path: Path, tool: str, available: bool) -> ToolResult:
    # Always use the final suffix — multi-dot filenames (e.g. 2606.28105v1.pdf)
    # would otherwise produce ".28105v1.pdf" which matches nothing in TOOL_SUPPORT.
    ext = path.suffix.lower()

    supported_exts = TOOL_SUPPORT.get(ext, set())
    supported = tool in supported_exts
    file_size_kb = round(path.stat().st_size / 1024, 1)

    # Infer file_type folder from parent dir name
    file_type = path.parent.name

    base = ToolResult(
        file=path.name,
        file_type=file_type,
        file_size_kb=file_size_kb,
        tool=tool,
        supported=supported,
        available=available,
        status="unsupported",
    )

    if not supported:
        return base
    if not available:
        return ToolResult(**{**asdict(base), "status": "unavailable"})

    # Skip very large files (>50 MB) to avoid OOM
    if path.stat().st_size > 50 * 1024 * 1024:
        return ToolResult(**{**asdict(base), "status": "skipped:too_large"})

    try:

        def _timeout_handler(signum, frame):
            raise TimeoutError(f"Exceeded {TIMEOUT_SECONDS}s")

        # Windows doesn't support SIGALRM; use threading timeout instead
        import threading
        result_holder: list = []
        error_holder: list = []

        def _run():
            try:
                text, elapsed = RUNNERS[tool](path)
                result_holder.append((text, elapsed))
            except Exception as e:
                error_holder.append(e)

        t = threading.Thread(target=_run, daemon=True)
        t.start()
        t.join(timeout=TIMEOUT_SECONDS)

        if t.is_alive():
            return ToolResult(**{**asdict(base), "status": "timeout",
                                 "error": f">{TIMEOUT_SECONDS}s"})
        if error_holder:
            raise error_holder[0]

        text, elapsed = result_holder[0]
        if not text or not text.strip():
            return ToolResult(**{**asdict(base), "status": "empty",
                                 "elapsed_seconds": round(elapsed, 3)})

        m = compute_metrics(tool, path.name, text, elapsed)
        return ToolResult(
            file=path.name,
            file_type=file_type,
            file_size_kb=file_size_kb,
            tool=tool,
            supported=supported,
            available=available,
            status="ok",
            token_count=m.token_count,
            char_count=m.char_count,
            heading_count=m.heading_count,
            table_count=m.table_count,
            noise_line_count=m.noise_line_count,
            duplicate_line_count=m.duplicate_line_count,
            avg_paragraph_tokens=m.avg_paragraph_tokens,
            quality_score=m.quality_score,
            elapsed_seconds=m.elapsed_seconds,
        )

    except Exception as e:
        return ToolResult(**{**asdict(base), "status": "error",
                             "error": str(e)[:200]})


def gather_files(corpus_dir: Path, type_filter: list[str] | None,
                 limit: int | None) -> list[Path]:
    files: list[Path] = []
    skip_kinds = {"youtube_url", "wikipedia_url"}

    for kind_dir in sorted(corpus_dir.iterdir()):
        if not kind_dir.is_dir():
            continue
        if kind_dir.name in skip_kinds:
            continue
        if type_filter and kind_dir.name not in type_filter:
            continue

        kind_files = sorted(kind_dir.iterdir())
        if limit:
            kind_files = kind_files[:limit]
        files.extend(kind_files)

    return files


def print_summary(results: list[ToolResult], available_tools: dict[str, bool]) -> None:
    from rich.console import Console
    from rich.table import Table

    console = Console()

    # ── availability table ────────────────────────────────────────────────────
    avail_table = Table(title="Tool Availability", show_header=True)
    avail_table.add_column("Tool", style="bold")
    avail_table.add_column("Installed", justify="center")
    for tool in ALL_TOOLS:
        sym = "[green]YES[/green]" if available_tools.get(tool) else "[red]NO[/red]"
        avail_table.add_row(tool, sym)
    console.print(avail_table)
    console.print()

    # ── per-type summary ──────────────────────────────────────────────────────
    from collections import defaultdict
    # {file_type: {tool: [ToolResult]}}
    by_type: dict[str, dict[str, list[ToolResult]]] = defaultdict(lambda: defaultdict(list))
    for r in results:
        by_type[r.file_type][r.tool].append(r)

    tools_to_show = [t for t in ALL_TOOLS if available_tools.get(t)]

    summary_table = Table(title="Results by File Type", show_header=True, show_lines=True)
    summary_table.add_column("Type", style="bold cyan", min_width=8)
    summary_table.add_column("Files", justify="right")
    for tool in tools_to_show:
        summary_table.add_column(f"{tool}\nok/err/uns", justify="center")
    summary_table.add_column("Token winner", justify="center")
    summary_table.add_column("Quality leader", justify="center")

    for ftype in sorted(by_type.keys()):
        tool_data = by_type[ftype]
        row = [ftype]

        total_files = max(len(v) for v in tool_data.values()) if tool_data else 0
        row.append(str(total_files))

        avg_tokens: dict[str, float] = {}
        for tool in tools_to_show:
            rlist = tool_data.get(tool, [])
            ok = sum(1 for r in rlist if r.status == "ok")
            err = sum(1 for r in rlist if r.status == "error")
            uns = sum(1 for r in rlist if r.status in ("unsupported", "unavailable"))
            color = "green" if ok > 0 else "dim"
            row.append(f"[{color}]{ok}[/{color}]/{err}/{uns}")
            tokens = [r.token_count for r in rlist if r.status == "ok" and r.token_count > 0]
            if tokens:
                avg_tokens[tool] = sum(tokens) / len(tokens)

        avg_quality: dict[str, float] = {}
        for tool in tools_to_show:
            rlist = tool_data.get(tool, [])
            scores = [r.quality_score for r in rlist if r.status == "ok"]
            if scores:
                avg_quality[tool] = sum(scores) / len(scores)

        if avg_tokens:
            winner = min(avg_tokens, key=avg_tokens.get)
            row.append(f"[bold green]{winner}[/bold green]")
        else:
            row.append("-")

        if avg_quality:
            qual_leader = max(avg_quality, key=avg_quality.get)
            best_q = avg_quality[qual_leader]
            q_color = "green" if best_q >= 80 else "yellow"
            row.append(f"[{q_color}]{qual_leader}({best_q:.0f})[/{q_color}]")
        else:
            row.append("-")

        summary_table.add_row(*row)

    console.print(summary_table)
    console.print()

    # ── overall stats ─────────────────────────────────────────────────────────
    overall = Table(title="Overall Stats (ok files only)", show_header=True)
    overall.add_column("Tool", style="bold")
    overall.add_column("Files OK", justify="right")
    overall.add_column("Avg tokens", justify="right")
    overall.add_column("Avg quality", justify="right")
    overall.add_column("Avg elapsed (s)", justify="right")
    overall.add_column("Noise lines avg", justify="right")
    overall.add_column("Types covered", justify="right")

    for tool in tools_to_show:
        ok_results = [r for r in results if r.tool == tool and r.status == "ok"]
        if not ok_results:
            overall.add_row(tool, "0", "-", "-", "-", "-", "-")
            continue
        avg_tok = sum(r.token_count for r in ok_results) / len(ok_results)
        avg_ela = sum(r.elapsed_seconds for r in ok_results) / len(ok_results)
        avg_noi = sum(r.noise_line_count for r in ok_results) / len(ok_results)
        avg_qual = sum(r.quality_score for r in ok_results) / len(ok_results)
        types_covered = len({r.file_type for r in ok_results})
        qual_color = "green" if avg_qual >= 80 else "yellow" if avg_qual >= 60 else "red"
        overall.add_row(
            tool,
            str(len(ok_results)),
            f"{avg_tok:,.0f}",
            f"[{qual_color}]{avg_qual:.0f}[/{qual_color}]",
            f"{avg_ela:.2f}",
            f"{avg_noi:.1f}",
            str(types_covered),
        )
    console.print(overall)
    console.print()

    # ── missing format coverage ───────────────────────────────────────────────
    missing = Table(title="Format coverage gaps (AksharaMD)", show_header=True)
    missing.add_column("Extension", style="bold red")
    missing.add_column("Currently handled by")
    for ext, tools in sorted(TOOL_SUPPORT.items()):
        if "aksharamd" not in tools and tools - {"aksharamd"}:
            others = ", ".join(sorted(tools - {"aksharamd"}))
            missing.add_row(ext, others)
    console.print(missing)

    # ── LLM input cost table ──────────────────────────────────────────────────
    LLM_PRICES = [
        ("Claude Sonnet 4",    3.00),
        ("GPT-4o",             2.50),
        ("Gemini 1.5 Flash",   0.075),
        ("Gemini 2.0 Flash",   0.10),
    ]
    cost_table = Table(
        title="Estimated LLM Input Cost per 1,000 Documents (tokens only — pending QA fidelity validation)",
        show_header=True
    )
    cost_table.add_column("Tool", style="bold")
    cost_table.add_column("Avg tokens", justify="right")
    for model, price in LLM_PRICES:
        cost_table.add_column(f"{model}\n(${price:.3f}/1M)", justify="right")

    for tool in tools_to_show:
        ok_results = [r for r in results if r.tool == tool and r.status == "ok"]
        if not ok_results:
            continue
        avg_tok = sum(r.token_count for r in ok_results) / len(ok_results)
        row = [tool, f"{avg_tok:,.0f}"]
        for model, price in LLM_PRICES:
            cost_per_1k = (avg_tok / 1_000_000) * price * 1000
            row.append(f"${cost_per_1k:.3f}")
        cost_table.add_row(*row)
    console.print(cost_table)
    console.print("[dim]Note: Cost estimates assume token count maps directly to LLM input tokens. "
                  "QA fidelity test required to validate information preservation.[/dim]")


def main() -> int:
    p = argparse.ArgumentParser(description="Multi-tool corpus benchmark")
    p.add_argument("--corpus-dir", default="../Downloads/benchmark_corpus",
                   help="Root of the downloaded benchmark corpus")
    p.add_argument("--output-dir", default="benchmark_results/corpus",
                   help="Where to write reports")
    p.add_argument("--tools", nargs="*", default=None,
                   help=f"Tools to run (default: all). Choices: {ALL_TOOLS}")
    p.add_argument("--types", nargs="*", default=None,
                   help="Corpus subdirs to include (e.g. pdf docx)")
    p.add_argument("--limit", type=int, default=None,
                   help="Max files per type (default: all)")
    p.add_argument("--no-rich", action="store_true",
                   help="Plain text output (no Rich tables)")
    args = p.parse_args()

    corpus_dir = Path(args.corpus_dir)
    if not corpus_dir.exists():
        print(f"Corpus not found: {corpus_dir}", file=sys.stderr)
        return 1

    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    import subprocess
    import sys as _sys
    def _get_version(pkg):
        try:
            r = subprocess.run([_sys.executable, "-m", "pip", "show", pkg],
                               capture_output=True, text=True)
            for line in r.stdout.splitlines():
                if line.startswith("Version:"):
                    return line.split(":", 1)[1].strip()
        except Exception:
            pass
        return "unknown"

    tools_to_run = args.tools or ALL_TOOLS
    unknown = set(tools_to_run) - set(ALL_TOOLS)
    if unknown:
        print(f"Unknown tools: {sorted(unknown)}. Valid: {ALL_TOOLS}", file=sys.stderr)
        return 1

    tool_versions = {t: _get_version(t if t != "aksharamd" else "aksharamd")
                     for t in tools_to_run}
    tool_versions["naive"] = _get_version("pymupdf")

    # Check availability
    available: dict[str, bool] = {}
    for tool in tools_to_run:
        available[tool] = check_tool_available(tool)
        status = "OK" if available[tool] else "MISSING (not installed)"
        print(f"  {tool:15} {status}")
    print()

    files = gather_files(corpus_dir, args.types, args.limit)
    print(f"Found {len(files)} files across {len({f.parent.name for f in files})} types.\n")

    results: list[ToolResult] = []
    total = len(files) * len(tools_to_run)
    done = 0

    for path in files:
        for tool in tools_to_run:
            done += 1
            ext = path.suffix.lower()
            supported = tool in TOOL_SUPPORT.get(ext, set())
            if not supported:
                results.append(ToolResult(
                    file=path.name, file_type=path.parent.name,
                    file_size_kb=round(path.stat().st_size / 1024, 1),
                    tool=tool, supported=False, available=available.get(tool, False),
                    status="unsupported",
                ))
                continue

            print(f"[{done:4}/{total}] {tool:14} {path.parent.name}/{path.name[:50]}", end=" ", flush=True)
            r = run_file(path, tool, available.get(tool, False))
            results.append(r)

            if r.status == "ok":
                print(f"ok  {r.token_count:>8,} tok  {r.elapsed_seconds:.2f}s")
            else:
                print(f"{r.status}  {r.error[:60]}")

    # ── write CSV ─────────────────────────────────────────────────────────────
    csv_path = out_dir / "results.csv"
    fields = list(ToolResult.__dataclass_fields__.keys())
    with csv_path.open("w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=fields)
        w.writeheader()
        for r in results:
            w.writerow(asdict(r))
    print(f"\nCSV  -> {csv_path}")

    # ── write JSON ────────────────────────────────────────────────────────────
    json_path = out_dir / "results.json"
    output = {
        "meta": {
            "date": __import__("datetime").datetime.utcnow().isoformat(),
            "corpus_dir": str(corpus_dir),
            "tool_versions": tool_versions,
            "tokenizer": "cl100k_base",
        },
        "results": [asdict(r) for r in results],
    }
    json_path.write_text(
        json.dumps(output, indent=2), encoding="utf-8"
    )
    print(f"JSON -> {json_path}")

    # ── summary ───────────────────────────────────────────────────────────────
    if not args.no_rich:
        try:
            print_summary(results, available)
        except ImportError:
            print("(install rich for pretty tables: pip install rich)")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
