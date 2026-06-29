#!/usr/bin/env python3
"""
Automated benchmark corpus builder.

Downloads or generates 20-50 real-world files per supported format so the
corpus_benchmark.py has statistically meaningful input.

Sources:
  pdf   - arXiv papers (CC-BY)
  html  - Wikipedia articles (CC-BY-SA)
  txt   - Project Gutenberg public-domain books
  epub  - Project Gutenberg public-domain books
  docx  - Synthetic (python-docx)
  pptx  - Synthetic (python-pptx)
  xlsx  - Synthetic (openpyxl)
  csv   - Synthetic tabular data
  json  - Synthetic structured data
  jsonl - Synthetic log / record data
  xml   - Synthetic config / report data
  rss   - Downloaded public RSS feeds
  atom  - Downloaded public ATOM feeds
  eml   - Synthetic email messages
  ipynb - Synthetic Jupyter notebooks
  jpg   - Synthetic PIL images with text (for OCR testing)
  png   - Synthetic PIL images with text (for OCR testing)
  zip   - Synthetic archives of text files

Usage:
  python -m benchmarks.corpus_builder --output-dir ../Downloads/benchmark_corpus
  python -m benchmarks.corpus_builder --output-dir corpus --types pdf docx html --count 25
  python -m benchmarks.corpus_builder --list-types
"""
from __future__ import annotations

import argparse
import io
import json
import os
import random
import sys
import time
import zipfile
from pathlib import Path
from typing import Callable

# ── constants ──────────────────────────────────────────────────────────────────

DEFAULT_COUNT = 25

ARXIV_CATEGORIES = [
    "cs.AI", "cs.CL", "cs.CV", "cs.LG", "cs.SE",
    "econ.GN", "q-bio.GN", "math.ST", "physics.data-an",
]

WIKIPEDIA_ARTICLES = [
    "Machine_learning", "Quantum_computing", "CRISPR", "Black_hole", "Climate_change",
    "World_War_II", "Roman_Empire", "Renaissance", "Industrial_Revolution", "Cold_War",
    "DNA", "Photosynthesis", "Evolution", "Neuroscience", "Cancer",
    "Internet", "Artificial_intelligence", "Blockchain", "Operating_system",
    "Python_(programming_language)", "Periodic_table", "Philosophy", "Democracy",
    "Classical_music", "William_Shakespeare", "Ancient_Egypt", "Moon_landing",
    "Amazon_rainforest", "Human_rights", "Supply_chain", "Transistor",
    "General_relativity", "Plate_tectonics", "Vaccination", "Genome",
    "Byzantine_Empire", "Silk_Road", "French_Revolution", "Mesopotamia",
    "Solar_energy", "Nuclear_fusion", "Stock_market", "Inflation",
    "Cognitive_science", "Linguistics", "Archaeology", "Sociology",
    "Urban_planning", "Oceanography",
]

GUTENBERG_IDS = [
    1342, 84, 2701, 1661, 74, 1232, 5200, 174, 345,
    768, 98, 219, 844, 514, 11, 25344, 2591, 2814,
    4300, 2600, 160, 541, 1400, 2554, 1260,
]

RSS_FEEDS = [
    "https://export.arxiv.org/rss/cs.AI",
    "https://export.arxiv.org/rss/cs.CL",
    "https://export.arxiv.org/rss/econ.GN",
    "https://www.nasa.gov/rss/dyn/breaking_news.rss",
    "https://feeds.bbci.co.uk/news/science_and_environment/rss.xml",
    "https://feeds.bbci.co.uk/news/technology/rss.xml",
    "https://feeds.bbci.co.uk/news/world/rss.xml",
    "https://rss.sciencedaily.com/all.xml",
    "https://www.sciencenews.org/feed",
    "https://physicsworld.com/feed/",
    "https://spectrum.ieee.org/rss/blog/tech-talk/fulltext",
    "https://www.nature.com/nature.rss",
    "https://feeds.aps.org/rss/recent/prl.rss",
    "https://api.quantamagazine.org/feed/",
    "https://www.newscientist.com/feed/home",
]

ATOM_FEEDS = [
    "https://export.arxiv.org/rss/cs.CV",
    "https://export.arxiv.org/rss/math.ST",
    "https://export.arxiv.org/rss/q-bio.GN",
    "https://export.arxiv.org/rss/physics.data-an",
    "https://www.w3.org/blog/news/feed/atom",
    "https://github.com/python/cpython/releases.atom",
    "https://github.com/microsoft/vscode/releases.atom",
    "https://github.com/openai/openai-python/releases.atom",
    "https://github.com/anthropics/anthropic-sdk-python/releases.atom",
    "https://blog.python.org/feeds/posts/default",
    "https://planet.scipy.org/atom.xml",
    "https://lwn.net/headlines/newrss",
    "https://www.phoronix.com/rss.php",
    "https://feeds.feedburner.com/oreilly/radar/atom",
    "https://sachachua.com/blog/feed/atom",
]

# ── helpers ────────────────────────────────────────────────────────────────────

def _get(url: str, timeout: int = 30) -> bytes | None:
    try:
        import requests
        r = requests.get(url, timeout=timeout, headers={
            "User-Agent": "AksharaMD-BenchmarkBuilder/1.0 (research; contact: benchmark@aksharamd.dev)"
        })
        if r.status_code == 200:
            return r.content
        print(f"    HTTP {r.status_code}: {url[:80]}")
        return None
    except Exception as e:
        print(f"    Error fetching {url[:80]}: {e}")
        return None


def _write(path: Path, data: bytes) -> bool:
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(data)
        return True
    except Exception as e:
        print(f"    Write error {path}: {e}")
        return False


def _slug(s: str, maxlen: int = 40) -> str:
    import re
    s = re.sub(r"[^\w\-]", "_", s.lower())
    s = re.sub(r"_+", "_", s).strip("_")
    return s[:maxlen]


# ── download: PDF from arXiv ──────────────────────────────────────────────────

def _download_pdf(out_dir: Path, count: int, overwrite: bool) -> int:
    import xml.etree.ElementTree as ET

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    ids_seen: set[str] = set()

    for cat in ARXIV_CATEGORIES:
        if saved >= count:
            break

        url = (
            f"https://export.arxiv.org/api/query"
            f"?search_query=cat:{cat}&max_results=15&sortBy=submittedDate&sortOrder=descending"
        )
        xml_bytes = _get(url)
        if not xml_bytes:
            continue

        try:
            root = ET.fromstring(xml_bytes)
            ns = {"atom": "http://www.w3.org/2005/Atom"}
            for entry in root.findall("atom:entry", ns):
                if saved >= count:
                    break
                id_url = (entry.findtext("atom:id", "", ns) or "").strip()
                arxiv_id = id_url.split("/abs/")[-1].replace("/", "_")
                if arxiv_id in ids_seen or not arxiv_id:
                    continue
                ids_seen.add(arxiv_id)

                dest = out_dir / f"{arxiv_id}.pdf"
                if dest.exists() and not overwrite:
                    saved += 1
                    continue

                pdf_url = id_url.replace("/abs/", "/pdf/")
                print(f"  PDF [{saved+1}/{count}] {arxiv_id}")
                data = _get(pdf_url, timeout=60)
                if data and data[:4] == b"%PDF":
                    _write(dest, data)
                    saved += 1
                time.sleep(0.5)
        except ET.ParseError as e:
            print(f"    XML parse error for {cat}: {e}")

    return saved


# ── download: HTML from Wikipedia ─────────────────────────────────────────────

def _download_html(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    articles = list(WIKIPEDIA_ARTICLES)
    random.shuffle(articles)
    saved = 0

    for title in articles:
        if saved >= count:
            break
        dest = out_dir / f"{_slug(title)}.html"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        url = f"https://en.wikipedia.org/api/rest_v1/page/html/{title}"
        print(f"  HTML [{saved+1}/{count}] {title}")
        data = _get(url)
        if data:
            _write(dest, data)
            saved += 1
        time.sleep(0.3)

    return saved


# ── download: TXT from Project Gutenberg ──────────────────────────────────────

def _download_txt(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    ids = list(GUTENBERG_IDS[:count])
    saved = 0

    for gid in ids:
        if saved >= count:
            break
        dest = out_dir / f"gutenberg_{gid}.txt"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        print(f"  TXT [{saved+1}/{count}] Gutenberg #{gid}")
        # Try multiple URL patterns
        for url in [
            f"https://www.gutenberg.org/cache/epub/{gid}/pg{gid}.txt",
            f"https://www.gutenberg.org/files/{gid}/{gid}-0.txt",
            f"https://www.gutenberg.org/files/{gid}/{gid}.txt",
        ]:
            data = _get(url, timeout=60)
            if data and len(data) > 1000:
                _write(dest, data)
                saved += 1
                break
        time.sleep(0.5)

    return saved


# ── download: EPUB from Project Gutenberg ────────────────────────────────────

def _download_epub(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    ids = list(GUTENBERG_IDS[:count])
    saved = 0

    for gid in ids:
        if saved >= count:
            break
        dest = out_dir / f"gutenberg_{gid}.epub"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        print(f"  EPUB [{saved+1}/{count}] Gutenberg #{gid}")
        for url in [
            f"https://www.gutenberg.org/ebooks/{gid}.epub.noimages",
            f"https://www.gutenberg.org/cache/epub/{gid}/pg{gid}.epub",
        ]:
            data = _get(url, timeout=60)
            if data and data[:2] == b"PK":   # EPUB is a ZIP
                _write(dest, data)
                saved += 1
                break
        time.sleep(0.5)

    return saved


# ── download: RSS feeds ────────────────────────────────────────────────────────

def _download_rss(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    feeds = list(RSS_FEEDS)
    random.shuffle(feeds)
    saved = 0

    for url in feeds:
        if saved >= count:
            break
        name = _slug(url.split("//")[-1].split("/")[0]) + f"_{saved+1}"
        dest = out_dir / f"{name}.rss"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        print(f"  RSS [{saved+1}/{count}] {url[:60]}")
        data = _get(url)
        if data and (b"<rss" in data[:500] or b"<channel" in data[:1000]):
            _write(dest, data)
            saved += 1
        elif data and b"<feed" in data[:500]:
            # Some feeds return Atom; save anyway
            _write(dest, data)
            saved += 1
        time.sleep(0.3)

    return saved


# ── download: ATOM feeds ──────────────────────────────────────────────────────

def _download_atom(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    feeds = list(ATOM_FEEDS)
    random.shuffle(feeds)
    saved = 0

    for url in feeds:
        if saved >= count:
            break
        name = _slug(url.split("//")[-1].split("/")[0]) + f"_{saved+1}"
        dest = out_dir / f"{name}.atom"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        print(f"  ATOM [{saved+1}/{count}] {url[:60]}")
        data = _get(url)
        if data and (b"<feed" in data[:500] or b"<rss" in data[:500]):
            _write(dest, data)
            saved += 1
        time.sleep(0.3)

    return saved


# ── generate: DOCX ────────────────────────────────────────────────────────────

_DOCX_TOPICS = [
    ("Quarterly Business Review", "business"),
    ("Technical Architecture Specification", "technical"),
    ("Market Research Analysis", "research"),
    ("Project Implementation Plan", "project"),
    ("Employee Handbook Overview", "policy"),
    ("Risk Assessment Report", "report"),
    ("Product Roadmap 2025", "roadmap"),
    ("Annual Financial Summary", "financial"),
    ("Customer Onboarding Guide", "guide"),
    ("API Integration Specification", "technical"),
    ("Competitive Analysis Report", "research"),
    ("Change Management Plan", "project"),
    ("Security Policy Document", "policy"),
    ("Data Governance Framework", "report"),
    ("Strategic Partnership Proposal", "business"),
    ("Software Requirements Specification", "technical"),
    ("Operations Manual", "guide"),
    ("Budget Allocation Review", "financial"),
    ("Performance Metrics Report", "report"),
    ("Training Program Curriculum", "policy"),
    ("Infrastructure Upgrade Plan", "technical"),
    ("Brand Guidelines Document", "guide"),
    ("Incident Response Playbook", "policy"),
    ("Research Methodology Overview", "research"),
    ("Stakeholder Communication Plan", "project"),
]

_LOREM = (
    "Organizations that invest in systematic documentation practices consistently "
    "outperform those that rely on institutional knowledge alone. The ability to "
    "capture, structure, and retrieve information efficiently is a core competency "
    "in modern enterprises. This section outlines the key considerations and "
    "recommended approaches for implementing effective documentation workflows."
)

_LOREM_SHORT = [
    "Effective communication is the foundation of any successful initiative.",
    "Data-driven decisions reduce uncertainty and improve outcome predictability.",
    "Stakeholder alignment must be established before execution begins.",
    "Regular reviews ensure that objectives remain relevant and achievable.",
    "Cross-functional collaboration accelerates delivery and reduces rework.",
    "Clear ownership and accountability are prerequisites for sustainable progress.",
    "Metrics should be leading indicators, not lagging measurements of failure.",
]


def _generate_docx(out_dir: Path, count: int, overwrite: bool) -> int:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        print("  SKIP docx: python-docx not installed")
        return 0

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(42)

    for i, (title, kind) in enumerate((_DOCX_TOPICS * 3)[:count]):
        dest = out_dir / f"doc_{i+1:02d}_{_slug(title)}.docx"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        doc = Document()
        doc.add_heading(title, level=1)
        doc.add_paragraph(f"Prepared by: Operations Team  |  Version {i+1}.0  |  Confidential")
        doc.add_heading("Executive Summary", level=2)
        doc.add_paragraph(_LOREM)
        doc.add_paragraph(rng.choice(_LOREM_SHORT))

        if kind in ("business", "financial", "report"):
            doc.add_heading("Key Metrics", level=2)
            table = doc.add_table(rows=1, cols=4)
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            for j, h in enumerate(["Metric", "Target", "Actual", "Variance"]):
                hdr[j].text = h
            for metric in ["Revenue", "Cost", "Margin", "NPS Score", "Retention"]:
                row = table.add_row().cells
                row[0].text = metric
                row[1].text = f"${rng.randint(100, 999):,}K"
                row[2].text = f"${rng.randint(80, 950):,}K"
                pct = rng.randint(-15, 20)
                row[3].text = f"{'+' if pct >= 0 else ''}{pct}%"

        elif kind in ("technical", "guide"):
            doc.add_heading("Architecture Overview", level=2)
            doc.add_paragraph(
                "The system is composed of three primary layers: the ingestion layer, "
                "the processing pipeline, and the delivery tier. Each layer is independently "
                "scalable and communicates via well-defined interfaces."
            )
            doc.add_heading("Component Summary", level=3)
            table = doc.add_table(rows=1, cols=3)
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            for j, h in enumerate(["Component", "Technology", "SLA"]):
                hdr[j].text = h
            components = [
                ("API Gateway", "nginx / Kong", "99.9%"),
                ("Auth Service", "OAuth 2.0 / JWT", "99.95%"),
                ("Data Pipeline", "Apache Kafka", "99.5%"),
                ("Storage Layer", "PostgreSQL + S3", "99.99%"),
                ("Cache Tier", "Redis Cluster", "99.9%"),
            ]
            for comp in components:
                row = table.add_row().cells
                for j, val in enumerate(comp):
                    row[j].text = val

        elif kind in ("project", "roadmap"):
            doc.add_heading("Project Timeline", level=2)
            table = doc.add_table(rows=1, cols=4)
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            for j, h in enumerate(["Phase", "Duration", "Owner", "Status"]):
                hdr[j].text = h
            phases = [
                ("Discovery", "2 weeks", "Product", "Complete"),
                ("Design", "3 weeks", "UX/Design", "In Progress"),
                ("Development", "8 weeks", "Engineering", "Planned"),
                ("Testing", "3 weeks", "QA", "Planned"),
                ("Deployment", "1 week", "DevOps", "Planned"),
            ]
            for phase in phases:
                row = table.add_row().cells
                for j, val in enumerate(phase):
                    row[j].text = val

        doc.add_heading("Recommendations", level=2)
        for bullet in rng.sample(_LOREM_SHORT, 4):
            doc.add_paragraph(bullet, style="List Bullet")

        doc.add_heading("Next Steps", level=2)
        doc.add_paragraph(
            "The team should convene within the next five business days to review these "
            "findings, assign owners to each recommendation, and establish measurable "
            "success criteria. Progress will be reviewed at the next steering committee meeting."
        )

        doc.save(str(dest))
        saved += 1
        print(f"  DOCX [{saved}/{count}] {dest.name}")

    return saved


# ── generate: PPTX ────────────────────────────────────────────────────────────

_PPTX_TOPICS = [
    "Q4 Business Results", "Product Vision 2025", "Technical Architecture Review",
    "Market Expansion Strategy", "Customer Success Overview", "Engineering All-Hands",
    "Investor Update", "Data Platform Roadmap", "Security Posture Review",
    "Team OKRs Q1", "Partnership Proposal", "ML Infrastructure Overview",
    "Go-to-Market Strategy", "Cost Optimization Plan", "Hiring Plan 2025",
    "Competitive Landscape", "API Strategy", "Cloud Migration Update",
    "User Research Findings", "Brand Refresh Proposal",
]


def _generate_pptx(out_dir: Path, count: int, overwrite: bool) -> int:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("  SKIP pptx: python-pptx not installed")
        return 0

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(99)

    for i, topic in enumerate((_PPTX_TOPICS * 3)[:count]):
        dest = out_dir / f"ppt_{i+1:02d}_{_slug(topic)}.pptx"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        prs = Presentation()
        slide_layout = prs.slide_layouts

        # Title slide
        slide = prs.slides.add_slide(slide_layout[0])
        slide.shapes.title.text = topic
        if slide.placeholders[1]:
            slide.placeholders[1].text = f"Presenter: Team Lead  |  {2024 + (i % 2)}"

        # Agenda
        slide = prs.slides.add_slide(slide_layout[1])
        slide.shapes.title.text = "Agenda"
        tf = slide.placeholders[1].text_frame
        tf.text = "Overview"
        for item in ["Key Metrics", "Analysis", "Recommendations", "Next Steps"]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1

        # Key findings slide
        for section in ["Key Findings", "Analysis", "Recommendations"]:
            slide = prs.slides.add_slide(slide_layout[1])
            slide.shapes.title.text = section
            tf = slide.placeholders[1].text_frame
            tf.text = rng.choice(_LOREM_SHORT)
            for bullet in rng.sample(_LOREM_SHORT, 3):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1

        # Metrics slide
        slide = prs.slides.add_slide(slide_layout[5])
        slide.shapes.title.text = "Key Metrics"

        prs.save(str(dest))
        saved += 1
        print(f"  PPTX [{saved}/{count}] {dest.name}")

    return saved


# ── generate: XLSX ────────────────────────────────────────────────────────────

def _generate_xlsx(out_dir: Path, count: int, overwrite: bool) -> int:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        print("  SKIP xlsx: openpyxl not installed")
        return 0

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(7)

    templates = [
        ("sales", _gen_xlsx_sales),
        ("budget", _gen_xlsx_budget),
        ("employees", _gen_xlsx_employees),
        ("inventory", _gen_xlsx_inventory),
        ("metrics", _gen_xlsx_metrics),
    ]

    for i in range(count):
        name, gen_fn = templates[i % len(templates)]
        dest = out_dir / f"xlsx_{i+1:02d}_{name}.xlsx"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        wb = gen_fn(rng, i)
        wb.save(str(dest))
        saved += 1
        print(f"  XLSX [{saved}/{count}] {dest.name}")

    return saved


def _gen_xlsx_sales(rng: random.Random, seed: int):
    import openpyxl
    from openpyxl.styles import Font, PatternFill
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sales Data"
    headers = ["Date", "Product", "Region", "Units", "Unit Price", "Revenue", "Target", "Attainment"]
    ws.append(headers)
    products = ["Widget A", "Widget B", "Service Pro", "Enterprise Suite", "Add-on Pack"]
    regions = ["North", "South", "East", "West", "International"]
    for row in range(40 + seed * 3):
        month = f"2024-{(row % 12)+1:02d}-01"
        product = rng.choice(products)
        region = rng.choice(regions)
        units = rng.randint(10, 500)
        price = rng.randint(50, 2000)
        revenue = units * price
        target = int(revenue * rng.uniform(0.8, 1.3))
        att = f"{revenue/target*100:.1f}%"
        ws.append([month, product, region, units, price, revenue, target, att])
    # Summary sheet
    ws2 = wb.create_sheet("Summary")
    ws2.append(["Region", "Total Revenue", "Avg Units"])
    for r in regions:
        ws2.append([r, rng.randint(50000, 500000), rng.randint(100, 1000)])
    return wb


def _gen_xlsx_budget(rng: random.Random, seed: int):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget vs Actual"
    ws.append(["Category", "Department", "Q1 Budget", "Q1 Actual", "Q2 Budget", "Q2 Actual", "YTD Variance"])
    categories = ["Headcount", "Software", "Infrastructure", "Marketing", "Travel", "Training", "Facilities"]
    depts = ["Engineering", "Sales", "Marketing", "HR", "Finance", "Operations"]
    for cat in categories:
        for dept in rng.sample(depts, 3):
            budget = rng.randint(10000, 200000)
            actual = int(budget * rng.uniform(0.7, 1.2))
            ws.append([cat, dept, budget, actual, int(budget*1.05), int(actual*rng.uniform(0.8, 1.1)),
                        f"{(actual-budget)/budget*100:.1f}%"])
    return wb


def _gen_xlsx_employees(rng: random.Random, seed: int):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Employees"
    ws.append(["ID", "Name", "Department", "Title", "Salary", "Hire Date", "Location", "Manager"])
    first = ["Alex", "Jordan", "Morgan", "Taylor", "Casey", "Riley", "Avery", "Quinn", "Drew", "Sage"]
    last = ["Smith", "Johnson", "Williams", "Brown", "Jones", "Garcia", "Miller", "Davis", "Wilson", "Moore"]
    depts = ["Engineering", "Sales", "Marketing", "HR", "Finance", "Product", "Operations", "Legal"]
    for i in range(50 + seed * 2):
        name = f"{rng.choice(first)} {rng.choice(last)}"
        dept = rng.choice(depts)
        salary = rng.randint(60000, 220000)
        month = rng.randint(1, 12)
        year = rng.randint(2018, 2024)
        ws.append([f"EMP{1000+i}", name, dept, f"Senior {dept} Specialist",
                    salary, f"{year}-{month:02d}-01", rng.choice(["NYC", "SF", "Austin", "Remote"]),
                    f"MGR{1000+rng.randint(0, 10)}"])
    return wb


def _gen_xlsx_inventory(rng: random.Random, seed: int):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Inventory"
    ws.append(["SKU", "Product Name", "Category", "Stock", "Reorder Point", "Unit Cost", "Supplier", "Last Updated"])
    categories = ["Electronics", "Software", "Hardware", "Accessories", "Services"]
    suppliers = ["Acme Corp", "Global Supply Co", "Tech Parts Inc", "Premium Goods Ltd"]
    for i in range(60 + seed):
        ws.append([
            f"SKU-{10000+i}", f"Product {chr(65+i%26)}-{i}", rng.choice(categories),
            rng.randint(0, 1000), rng.randint(10, 100), round(rng.uniform(5, 500), 2),
            rng.choice(suppliers), f"2024-{rng.randint(1,12):02d}-{rng.randint(1,28):02d}"
        ])
    return wb


def _gen_xlsx_metrics(rng: random.Random, seed: int):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "KPI Dashboard"
    ws.append(["Month", "DAU", "MAU", "Revenue", "Churn Rate", "NPS", "Support Tickets", "CSAT"])
    for month in range(1, 25):
        year = 2023 + (month > 12)
        m = (month - 1) % 12 + 1
        ws.append([
            f"{year}-{m:02d}", rng.randint(10000, 100000), rng.randint(50000, 500000),
            rng.randint(100000, 2000000), round(rng.uniform(0.5, 8.0), 2),
            rng.randint(20, 80), rng.randint(100, 2000), round(rng.uniform(3.0, 5.0), 1)
        ])
    return wb


# ── generate: CSV ─────────────────────────────────────────────────────────────

_CSV_SCHEMAS = [
    ("transactions", ["tx_id", "date", "account", "category", "amount", "currency", "merchant", "status"]),
    ("web_logs", ["timestamp", "ip", "method", "path", "status_code", "bytes", "referrer", "user_agent"]),
    ("sensors", ["sensor_id", "timestamp", "temperature", "humidity", "pressure", "location", "alert"]),
    ("products", ["id", "name", "category", "price", "stock", "rating", "reviews", "supplier"]),
    ("events", ["event_id", "user_id", "event_type", "page", "session", "timestamp", "duration_ms", "country"]),
]


def _generate_csv(out_dir: Path, count: int, overwrite: bool) -> int:
    import csv as csv_mod

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(13)

    for i in range(count):
        schema_name, headers = _CSV_SCHEMAS[i % len(_CSV_SCHEMAS)]
        dest = out_dir / f"csv_{i+1:02d}_{schema_name}.csv"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        rows_count = 80 + i * 5
        with dest.open("w", newline="", encoding="utf-8") as f:
            w = csv_mod.writer(f)
            w.writerow(headers)
            for r in range(rows_count):
                if schema_name == "transactions":
                    w.writerow([
                        f"TX{100000+r}", f"2024-{rng.randint(1,12):02d}-{rng.randint(1,28):02d}",
                        f"ACC-{rng.randint(1000, 9999)}", rng.choice(["Food", "Travel", "Software", "Hardware", "Services"]),
                        round(rng.uniform(-5000, 10000), 2), rng.choice(["USD", "EUR", "GBP"]),
                        rng.choice(["Amazon", "Google", "Stripe", "Acme", "Local Store"]),
                        rng.choice(["completed", "pending", "failed"])
                    ])
                elif schema_name == "web_logs":
                    w.writerow([
                        f"2024-01-{rng.randint(1,28):02d}T{rng.randint(0,23):02d}:{rng.randint(0,59):02d}:00Z",
                        f"{rng.randint(1,255)}.{rng.randint(0,255)}.{rng.randint(0,255)}.{rng.randint(0,255)}",
                        rng.choice(["GET", "POST", "PUT", "DELETE"]),
                        rng.choice(["/api/users", "/api/data", "/health", "/v1/process", "/static/index.html"]),
                        rng.choice([200, 200, 200, 201, 400, 404, 500]),
                        rng.randint(100, 50000), "https://example.com", "Mozilla/5.0"
                    ])
                elif schema_name == "sensors":
                    w.writerow([
                        f"SENSOR-{rng.randint(1, 50):03d}",
                        f"2024-01-{rng.randint(1,28):02d}T{rng.randint(0,23):02d}:00:00Z",
                        round(rng.uniform(-10, 45), 2), round(rng.uniform(20, 95), 2),
                        round(rng.uniform(950, 1050), 2),
                        rng.choice(["Floor 1", "Floor 2", "Roof", "Basement", "Server Room"]),
                        rng.choice(["OK", "OK", "OK", "WARNING", "CRITICAL"])
                    ])
                else:
                    w.writerow([f"row_{r}"] + [rng.randint(1, 1000) for _ in headers[1:]])

        saved += 1
        print(f"  CSV [{saved}/{count}] {dest.name}")

    return saved


# ── generate: JSON ────────────────────────────────────────────────────────────

def _generate_json(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(17)

    templates = ["api_response", "config", "schema", "graph", "report"]

    for i in range(count):
        kind = templates[i % len(templates)]
        dest = out_dir / f"json_{i+1:02d}_{kind}.json"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        if kind == "api_response":
            data = {
                "status": "success", "code": 200,
                "meta": {"page": i+1, "per_page": 25, "total": 1000},
                "data": [
                    {"id": j, "name": f"Item {j}", "score": round(rng.uniform(0, 100), 2),
                     "tags": rng.sample(["a", "b", "c", "d", "e"], k=2),
                     "active": rng.choice([True, False]),
                     "created_at": f"2024-{rng.randint(1,12):02d}-{rng.randint(1,28):02d}"}
                    for j in range(50)
                ]
            }
        elif kind == "config":
            data = {
                "version": "2.0",
                "environment": rng.choice(["production", "staging", "development"]),
                "database": {
                    "host": "db.internal", "port": 5432, "name": f"app_db_{i}",
                    "pool_size": rng.randint(5, 50), "ssl": True,
                    "replicas": [f"replica-{k}.internal" for k in range(3)]
                },
                "cache": {"backend": "redis", "ttl": 3600, "max_entries": 100000},
                "features": {f"feature_{k}": rng.choice([True, False]) for k in range(20)},
                "limits": {"rate_limit": 1000, "max_payload_kb": 512, "timeout_seconds": 30},
            }
        elif kind == "schema":
            data = {
                "$schema": "https://json-schema.org/draft/2020-12",
                "title": f"Schema_{i}", "type": "object",
                "properties": {
                    f"field_{k}": {
                        "type": rng.choice(["string", "integer", "boolean", "array"]),
                        "description": f"Field {k} description"
                    }
                    for k in range(30)
                },
                "required": [f"field_{k}" for k in range(5)]
            }
        elif kind == "graph":
            nodes = [{"id": k, "label": f"Node_{k}", "weight": rng.randint(1, 100)} for k in range(30)]
            edges = [{"source": rng.randint(0, 29), "target": rng.randint(0, 29),
                      "weight": round(rng.uniform(0, 1), 3)} for _ in range(60)]
            data = {"nodes": nodes, "edges": edges, "directed": True}
        else:  # report
            data = {
                "report_id": f"RPT-{2024}{i:04d}", "generated_at": "2024-01-01T00:00:00Z",
                "summary": {"total_items": rng.randint(100, 10000), "processed": rng.randint(80, 9999),
                             "errors": rng.randint(0, 50)},
                "sections": [
                    {"title": f"Section {k}", "content": rng.choice(_LOREM_SHORT),
                     "metrics": {f"metric_{m}": rng.randint(1, 1000) for m in range(5)}}
                    for k in range(10)
                ]
            }

        dest.write_text(json.dumps(data, indent=2), encoding="utf-8")
        saved += 1
        print(f"  JSON [{saved}/{count}] {dest.name}")

    return saved


# ── generate: JSONL ───────────────────────────────────────────────────────────

def _generate_jsonl(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(23)

    kinds = ["logs", "events", "records", "stream"]

    for i in range(count):
        kind = kinds[i % len(kinds)]
        dest = out_dir / f"jsonl_{i+1:02d}_{kind}.jsonl"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        lines = []
        for r in range(200 + i * 10):
            if kind == "logs":
                lines.append(json.dumps({
                    "ts": f"2024-01-{rng.randint(1,28):02d}T{rng.randint(0,23):02d}:{rng.randint(0,59):02d}:{rng.randint(0,59):02d}Z",
                    "level": rng.choice(["INFO", "INFO", "INFO", "WARN", "ERROR"]),
                    "service": rng.choice(["api", "worker", "scheduler", "auth", "storage"]),
                    "msg": rng.choice(["Request processed", "Task completed", "Cache miss", "Connection retry", "Rate limit hit"]),
                    "duration_ms": rng.randint(1, 5000), "trace_id": f"{rng.randint(10**15, 10**16-1):016x}"
                }))
            elif kind == "events":
                lines.append(json.dumps({
                    "event": rng.choice(["page_view", "click", "signup", "purchase", "logout"]),
                    "user_id": f"u{rng.randint(1000, 9999)}",
                    "session": f"s{rng.randint(100000, 999999)}",
                    "properties": {"page": rng.choice(["/home", "/docs", "/pricing", "/login"]),
                                   "referrer": rng.choice(["google", "direct", "email"])}
                }))
            else:
                lines.append(json.dumps({"id": r, "value": round(rng.gauss(100, 15), 4),
                                          "category": rng.choice(["A", "B", "C"]),
                                          "valid": rng.choice([True, False])}))

        dest.write_text("\n".join(lines), encoding="utf-8")
        saved += 1
        print(f"  JSONL [{saved}/{count}] {dest.name}")

    return saved


# ── generate: XML ─────────────────────────────────────────────────────────────

def _generate_xml(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(31)

    for i in range(count):
        dest = out_dir / f"xml_{i+1:02d}.xml"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        kind = ["catalog", "report", "config"][i % 3]
        lines = ['<?xml version="1.0" encoding="UTF-8"?>']

        if kind == "catalog":
            lines.append("<catalog>")
            for j in range(30):
                lines.append(f'  <item id="{j}" category="{rng.choice(["A","B","C"])}">')
                lines.append(f"    <name>Product {j}</name>")
                lines.append(f"    <price currency=\"USD\">{round(rng.uniform(10, 1000), 2)}</price>")
                lines.append(f"    <stock>{rng.randint(0, 500)}</stock>")
                lines.append(f"    <description>{rng.choice(_LOREM_SHORT)}</description>")
                lines.append("  </item>")
            lines.append("</catalog>")
        elif kind == "report":
            lines.append('<report xmlns="urn:aksharamd:report" version="1.0">')
            lines.append(f"  <title>Report {i}</title>")
            lines.append(f"  <generated>2024-01-{rng.randint(1,28):02d}</generated>")
            lines.append("  <sections>")
            for s in range(5):
                lines.append(f'    <section id="{s}">')
                lines.append(f"      <heading>Section {s+1}</heading>")
                lines.append(f"      <content>{rng.choice(_LOREM_SHORT)}</content>")
                lines.append("      <metrics>")
                for m in range(4):
                    lines.append(f'        <metric name="KPI_{m}" value="{rng.randint(1, 100)}" unit="pct"/>')
                lines.append("      </metrics>")
                lines.append("    </section>")
            lines.append("  </sections>")
            lines.append("</report>")
        else:
            lines.append("<configuration>")
            for key in ["database", "cache", "auth", "logging", "api"]:
                lines.append(f"  <{key}>")
                lines.append(f"    <host>{key}.internal.example.com</host>")
                lines.append(f"    <port>{rng.randint(1024, 65535)}</port>")
                lines.append(f"    <enabled>{str(rng.choice([True, False])).lower()}</enabled>")
                lines.append(f"  </{key}>")
            lines.append("</configuration>")

        dest.write_text("\n".join(lines), encoding="utf-8")
        saved += 1
        print(f"  XML [{saved}/{count}] {dest.name}")

    return saved


# ── generate: EML ─────────────────────────────────────────────────────────────

_EMAIL_SUBJECTS = [
    "Q4 Planning Meeting Follow-up",
    "RE: Partnership Agreement Review",
    "Action Required: Budget Approval",
    "Weekly Engineering Sync Notes",
    "Invitation: Product Demo - Thursday 2pm",
    "RE: Customer Escalation - Case #48291",
    "FWD: Security Patch Release Notes",
    "Project Kickoff - Next Steps",
    "RE: Interview Feedback - Candidate A",
    "Monthly KPI Report Attached",
    "RE: Contract Renewal Discussion",
    "Alert: System Maintenance Window",
    "Announcement: New Hire - Alex Chen",
    "RE: Legal Review Required",
    "Quarterly All-Hands Recording Available",
]


def _generate_eml(out_dir: Path, count: int, overwrite: bool) -> int:
    import email.mime.multipart
    import email.mime.text

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(37)
    names = ["alex.chen@example.com", "sarah.johnson@company.org", "mike.williams@corp.net",
             "team-lead@startup.io", "support@service.com", "noreply@platform.dev"]

    for i in range(count):
        dest = out_dir / f"email_{i+1:02d}.eml"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        subject = _EMAIL_SUBJECTS[i % len(_EMAIL_SUBJECTS)]
        sender = rng.choice(names)
        recipient = rng.choice([n for n in names if n != sender])

        body = f"""Dear Team,

{rng.choice(_LOREM_SHORT)}

{_LOREM}

Key points from our discussion:
1. {rng.choice(_LOREM_SHORT)}
2. {rng.choice(_LOREM_SHORT)}
3. {rng.choice(_LOREM_SHORT)}

Please review the attached documentation and provide feedback by end of week.

The next meeting is scheduled for Friday at 3pm EST. Please confirm your attendance.

Best regards,
{sender.split('@')[0].replace('.', ' ').title()}
{sender}
"""

        msg = email.mime.multipart.MIMEMultipart()
        msg["From"] = sender
        msg["To"] = recipient
        msg["Subject"] = subject
        msg["Date"] = f"Mon, {rng.randint(1,28):02d} Jan 2024 {rng.randint(8,17):02d}:30:00 +0000"
        msg["Message-ID"] = f"<{rng.randint(10**10, 10**11-1)}.{rng.randint(10**6, 10**7-1)}@example.com>"
        msg.attach(email.mime.text.MIMEText(body, "plain"))

        dest.write_bytes(msg.as_bytes())
        saved += 1
        print(f"  EML [{saved}/{count}] {dest.name}")

    return saved


# ── generate: IPYNB ───────────────────────────────────────────────────────────

_NOTEBOOK_TOPICS = [
    ("data_analysis", "Exploratory Data Analysis"),
    ("ml_training", "Machine Learning Model Training"),
    ("visualization", "Data Visualization"),
    ("text_processing", "Natural Language Processing"),
    ("statistics", "Statistical Analysis"),
]


def _generate_ipynb(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(41)

    for i in range(count):
        kind, title = _NOTEBOOK_TOPICS[i % len(_NOTEBOOK_TOPICS)]
        dest = out_dir / f"nb_{i+1:02d}_{kind}.ipynb"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        cells = []

        # Markdown header cell
        cells.append({
            "cell_type": "markdown", "metadata": {},
            "source": [f"# {title}\n\n", rng.choice(_LOREM_SHORT), "\n\n",
                        "## Setup\n\nImport required libraries and configure the environment."]
        })

        # Code: imports
        cells.append({
            "cell_type": "code", "metadata": {}, "execution_count": 1,
            "outputs": [],
            "source": ["import numpy as np\nimport pandas as pd\nimport matplotlib.pyplot as plt\n",
                        "from pathlib import Path\nimport json, time\n\n",
                        f"# {title} notebook - auto-generated for benchmarking\n",
                        "print('Environment ready')"]
        })

        # Markdown section
        cells.append({
            "cell_type": "markdown", "metadata": {},
            "source": ["## Data Loading\n\n", _LOREM]
        })

        # Code: data generation
        cells.append({
            "cell_type": "code", "metadata": {}, "execution_count": 2,
            "outputs": [{"name": "stdout", "output_type": "stream",
                          "text": ["Dataset shape: (1000, 5)\n"]}],
            "source": [
                "np.random.seed(42)\n",
                f"n_samples = {rng.randint(500, 5000)}\n",
                "data = pd.DataFrame({\n",
                "    'feature_1': np.random.normal(0, 1, n_samples),\n",
                "    'feature_2': np.random.uniform(0, 100, n_samples),\n",
                "    'category': np.random.choice(['A', 'B', 'C'], n_samples),\n",
                "    'label': np.random.randint(0, 2, n_samples),\n",
                "})\n",
                "print(f'Dataset shape: {data.shape}')\n",
                "data.head()"
            ]
        })

        # Markdown analysis
        cells.append({
            "cell_type": "markdown", "metadata": {},
            "source": ["## Analysis\n\n", rng.choice(_LOREM_SHORT), "\n\n",
                        "Key observations:\n- ", "\n- ".join(rng.sample(_LOREM_SHORT, 3))]
        })

        # Code: analysis
        cells.append({
            "cell_type": "code", "metadata": {}, "execution_count": 3,
            "outputs": [],
            "source": [
                "# Compute summary statistics\n",
                "summary = data.describe()\n",
                "print(summary)\n\n",
                "# Group by category\n",
                "grouped = data.groupby('category').mean(numeric_only=True)\n",
                "print(grouped)"
            ]
        })

        # Conclusion
        cells.append({
            "cell_type": "markdown", "metadata": {},
            "source": ["## Conclusion\n\n", rng.choice(_LOREM_SHORT), "\n\n",
                        "Next steps:\n1. ", "\n2. ".join(rng.sample(_LOREM_SHORT, 2))]
        })

        nb = {
            "nbformat": 4, "nbformat_minor": 5,
            "metadata": {"kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
                          "language_info": {"name": "python", "version": "3.11.0"}},
            "cells": cells,
        }

        dest.write_text(json.dumps(nb, indent=1), encoding="utf-8")
        saved += 1
        print(f"  IPYNB [{saved}/{count}] {dest.name}")

    return saved


# ── generate: JPG / PNG (synthetic OCR images) ───────────────────────────────

def _generate_images(out_dir: Path, count: int, overwrite: bool, ext: str) -> int:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print(f"  SKIP {ext}: Pillow not installed")
        return 0

    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(53)

    templates = [
        "invoice", "report_page", "form", "certificate", "memo"
    ]

    for i in range(count):
        kind = templates[i % len(templates)]
        dest = out_dir / f"img_{i+1:02d}_{kind}.{ext}"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        # Create a white A4-ish image (595 x 842 pixels at 72dpi)
        img = Image.new("RGB", (595, 842), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Try to get a readable font, fall back to default
        try:
            font_title = ImageFont.truetype("arial.ttf", 24)
            font_body = ImageFont.truetype("arial.ttf", 14)
            font_small = ImageFont.truetype("arial.ttf", 11)
        except Exception:
            font_title = ImageFont.load_default()
            font_body = font_title
            font_small = font_title

        if kind == "invoice":
            draw.text((40, 40), "INVOICE", font=font_title, fill=(0, 0, 0))
            draw.text((40, 80), f"Invoice #INV-{2024000+i}", font=font_body, fill=(50, 50, 50))
            draw.text((40, 100), f"Date: 2024-{rng.randint(1,12):02d}-{rng.randint(1,28):02d}", font=font_body, fill=(50, 50, 50))
            draw.text((40, 130), "Bill To:", font=font_body, fill=(0, 0, 0))
            draw.text((40, 150), "Acme Corporation", font=font_body, fill=(50, 50, 50))
            draw.text((40, 170), "123 Business Ave, New York, NY 10001", font=font_small, fill=(80, 80, 80))
            y = 220
            draw.text((40, y), "Description                    Qty    Unit Price    Total", font=font_small, fill=(0, 0, 0))
            y += 20
            draw.line([(40, y), (555, y)], fill=(200, 200, 200))
            y += 10
            for item in ["Professional Services", "Software License", "Support Plan", "Training"]:
                qty = rng.randint(1, 10)
                price = rng.randint(100, 2000)
                draw.text((40, y), f"{item:<30} {qty:>3}    ${price:>10,.2f}    ${qty*price:>10,.2f}", font=font_small, fill=(50, 50, 50))
                y += 20
            y += 20
            draw.text((40, y), f"Total Due: ${rng.randint(1000, 50000):,}.00", font=font_body, fill=(0, 0, 150))

        elif kind == "report_page":
            draw.text((40, 40), f"Monthly Performance Report", font=font_title, fill=(0, 0, 0))
            draw.text((40, 80), f"Period: {rng.choice(['January','February','March','April','May'])} 2024", font=font_body, fill=(80, 80, 80))
            y = 120
            draw.text((40, y), "Executive Summary", font=font_body, fill=(0, 0, 0))
            y += 25
            # Word-wrap the lorem text
            words = _LOREM.split()
            line, lines_out = [], []
            for w in words:
                if len(" ".join(line + [w])) < 65:
                    line.append(w)
                else:
                    lines_out.append(" ".join(line))
                    line = [w]
            if line:
                lines_out.append(" ".join(line))
            for ln in lines_out[:8]:
                draw.text((40, y), ln, font=font_small, fill=(60, 60, 60))
                y += 16

        else:
            draw.text((40, 40), kind.replace("_", " ").title(), font=font_title, fill=(0, 0, 0))
            y = 100
            for bullet in rng.sample(_LOREM_SHORT, 5):
                draw.text((40, y), f"- {bullet}", font=font_small, fill=(50, 50, 50))
                y += 24

        format_name = "JPEG" if ext == "jpg" else "PNG"
        img.save(str(dest), format_name, quality=85 if ext == "jpg" else None)
        saved += 1
        print(f"  {ext.upper()} [{saved}/{count}] {dest.name}")

    return saved


# ── generate: ZIP ─────────────────────────────────────────────────────────────

def _generate_zip(out_dir: Path, count: int, overwrite: bool) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    rng = random.Random(59)

    for i in range(count):
        dest = out_dir / f"archive_{i+1:02d}.zip"
        if dest.exists() and not overwrite:
            saved += 1
            continue

        n_files = rng.randint(5, 20)
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            # Add a README
            zf.writestr("README.txt", f"Archive #{i+1}\n\n{_LOREM}\n\n{rng.choice(_LOREM_SHORT)}")
            # Add some CSV data
            csv_content = "id,name,value,category\n"
            for r in range(50):
                csv_content += f"{r},item_{r},{rng.randint(1,1000)},{rng.choice(['A','B','C'])}\n"
            zf.writestr("data/records.csv", csv_content)
            # Add some JSON config
            config = {"version": "1.0", "settings": {f"key_{k}": rng.randint(1, 100) for k in range(10)}}
            zf.writestr("config/settings.json", json.dumps(config, indent=2))
            # Add multiple text files
            for j in range(min(n_files, 10)):
                content = f"Document {j+1}\n\n" + "\n\n".join(rng.sample(_LOREM_SHORT * 3, 4)) + "\n"
                zf.writestr(f"docs/doc_{j+1:02d}.txt", content)

        dest.write_bytes(buf.getvalue())
        saved += 1
        print(f"  ZIP [{saved}/{count}] {dest.name}")

    return saved


# ── registry ───────────────────────────────────────────────────────────────────

GeneratorFn = Callable[[Path, int, bool], int]

GENERATORS: dict[str, tuple[str, GeneratorFn]] = {
    "pdf":   ("Download from arXiv (CC-BY papers)",    lambda d, n, o: _download_pdf(d, n, o)),
    "html":  ("Download from Wikipedia (CC-BY-SA)",    lambda d, n, o: _download_html(d, n, o)),
    "txt":   ("Download from Project Gutenberg",       lambda d, n, o: _download_txt(d, n, o)),
    "epub":  ("Download from Project Gutenberg",       lambda d, n, o: _download_epub(d, n, o)),
    "rss":   ("Download public RSS feeds",             lambda d, n, o: _download_rss(d, n, o)),
    "atom":  ("Download public ATOM feeds",            lambda d, n, o: _download_atom(d, n, o)),
    "docx":  ("Synthetic (python-docx)",               lambda d, n, o: _generate_docx(d, n, o)),
    "pptx":  ("Synthetic (python-pptx)",               lambda d, n, o: _generate_pptx(d, n, o)),
    "xlsx":  ("Synthetic (openpyxl)",                  lambda d, n, o: _generate_xlsx(d, n, o)),
    "csv":   ("Synthetic tabular data",                lambda d, n, o: _generate_csv(d, n, o)),
    "json":  ("Synthetic structured JSON",             lambda d, n, o: _generate_json(d, n, o)),
    "jsonl": ("Synthetic JSONL logs/events",           lambda d, n, o: _generate_jsonl(d, n, o)),
    "xml":   ("Synthetic XML documents",               lambda d, n, o: _generate_xml(d, n, o)),
    "eml":   ("Synthetic email messages",              lambda d, n, o: _generate_eml(d, n, o)),
    "ipynb": ("Synthetic Jupyter notebooks",           lambda d, n, o: _generate_ipynb(d, n, o)),
    "jpg":   ("Synthetic OCR test images (Pillow)",    lambda d, n, o: _generate_images(d, n, o, "jpg")),
    "png":   ("Synthetic OCR test images (Pillow)",    lambda d, n, o: _generate_images(d, n, o, "png")),
    "zip":   ("Synthetic ZIP archives",                lambda d, n, o: _generate_zip(d, n, o)),
}


# ── CLI ────────────────────────────────────────────────────────────────────────

def main() -> int:
    p = argparse.ArgumentParser(
        description="Build benchmark corpus: download or generate files per format type.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    p.add_argument("--output-dir", default="../Downloads/benchmark_corpus",
                   help="Root directory for corpus (default: ../Downloads/benchmark_corpus)")
    p.add_argument("--types", nargs="*", default=None,
                   help=f"Types to build (default: all). Choices: {sorted(GENERATORS)}")
    p.add_argument("--count", type=int, default=DEFAULT_COUNT,
                   help=f"Target files per type (default: {DEFAULT_COUNT})")
    p.add_argument("--overwrite", action="store_true",
                   help="Re-download/regenerate even if files already exist")
    p.add_argument("--list-types", action="store_true",
                   help="List available types and exit")
    args = p.parse_args()

    if args.list_types:
        print(f"\nAvailable types ({len(GENERATORS)}):\n")
        for t, (desc, _) in sorted(GENERATORS.items()):
            print(f"  {t:<8} {desc}")
        print()
        return 0

    types = args.types or sorted(GENERATORS.keys())
    unknown = set(types) - set(GENERATORS.keys())
    if unknown:
        print(f"Unknown types: {sorted(unknown)}. Valid: {sorted(GENERATORS.keys())}", file=sys.stderr)
        return 1

    root = Path(args.output_dir)
    root.mkdir(parents=True, exist_ok=True)

    print(f"\nCorpus builder — target: {args.count} files/type  output: {root}\n")
    totals: dict[str, int] = {}

    for t in types:
        desc, fn = GENERATORS[t]
        out = root / t
        print(f"\n[{t.upper()}] {desc}")
        try:
            n = fn(out, args.count, args.overwrite)
            totals[t] = n
            print(f"  -> {n} files in {out}")
        except Exception as e:
            print(f"  ERROR: {e}")
            totals[t] = 0

    print(f"\n{'='*50}")
    print("Corpus build complete:\n")
    total_files = 0
    for t, n in sorted(totals.items()):
        status = "OK" if n >= args.count else f"PARTIAL ({n}/{args.count})"
        print(f"  {t:<8} {n:>3} files  {status}")
        total_files += n
    print(f"\n  Total: {total_files} files across {len(totals)} types")
    print(f"  Path:  {root.resolve()}\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
