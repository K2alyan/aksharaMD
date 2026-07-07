#!/usr/bin/env python3
"""Build the public reproducible benchmark corpus.

Downloads all 34 py-pdf/sample-files PDFs and generates 10 synthetic
variants per format (100 files total) into benchmarks/.public_corpus/.

Usage:
    python benchmarks/build_public_corpus.py [options]

Modes:
    (default / --full)   All 34 PDFs + 10 variants per format = 134 files
    --smoke              10 smoke PDFs + 1 variant per format = 20 files

Options:
    --dry-run            Print actions without writing files
    --skip-pdf           Skip PDF downloads entirely
    --max-download-mb M  Skip PDFs that would push cumulative download past M MB
                         (default: 100, 0 = unlimited)
"""
from __future__ import annotations

import argparse
import io
import json
import sys
import time
import urllib.error
import urllib.request
import zipfile
from pathlib import Path

REPO_ROOT = Path(__file__).parent.parent
BENCHMARKS = Path(__file__).parent
MANIFEST_PATH = BENCHMARKS / "public_corpus_manifest.json"
CORPUS_ROOT = BENCHMARKS / ".public_corpus"

_RETRY_ATTEMPTS = 3
_RETRY_DELAY = 2.0


def _load_manifest() -> dict:
    with open(MANIFEST_PATH, encoding="utf-8") as fh:
        return json.load(fh)


def _get_remote_size(url: str) -> int | None:
    """Return Content-Length in bytes via HEAD request, or None if unavailable."""
    try:
        req = urllib.request.Request(
            url, method="HEAD", headers={"User-Agent": "aksharamd-benchmark/1.0"}
        )
        with urllib.request.urlopen(req, timeout=10) as resp:
            length = resp.headers.get("Content-Length")
            return int(length) if length else None
    except Exception:
        return None


def _download(url: str, dest: Path, dry_run: bool = False) -> tuple[bool, int]:
    """Download url to dest. Returns (success, bytes_downloaded)."""
    dest.parent.mkdir(parents=True, exist_ok=True)
    if dest.exists():
        print(f"  skip (cached): {dest.name}")
        return True, 0
    if dry_run:
        print(f"  [dry-run] would download: {url}")
        return True, 0
    for attempt in range(1, _RETRY_ATTEMPTS + 1):
        try:
            req = urllib.request.Request(url, headers={"User-Agent": "aksharamd-benchmark/1.0"})
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = resp.read()
            dest.write_bytes(data)
            print(f"  downloaded ({len(data):,} B): {dest.name}")
            return True, len(data)
        except urllib.error.HTTPError as exc:
            print(f"  HTTP {exc.code} on attempt {attempt}/{_RETRY_ATTEMPTS}: {url}")
            if exc.code in (403, 404):
                return False, 0
        except Exception as exc:
            print(f"  error on attempt {attempt}/{_RETRY_ATTEMPTS}: {exc}")
        if attempt < _RETRY_ATTEMPTS:
            time.sleep(_RETRY_DELAY)
    return False, 0


# ── Synthetic builders ────────────────────────────────────────────────────────
# Each builder: _build_<fmt>(dest, variant=1) where variant is 1-indexed (1..10).
# Variant 1 is the "standard" fixture equivalent to the previous single-variant corpus.


def _build_docx(dest: Path, variant: int = 1) -> None:
    try:
        from docx import Document
        doc = Document()
        if variant == 1:  # headings-paragraphs
            doc.add_heading("AksharaMD Benchmark Document", 0)
            doc.add_paragraph(
                "This is a synthetic DOCX fixture for the AksharaMD public benchmark. "
                "It contains headings, paragraphs, and a table."
            )
            doc.add_heading("Section 1: Introduction", level=1)
            doc.add_paragraph(
                "The quick brown fox jumps over the lazy dog. "
                "Pack my box with five dozen liquor jugs."
            )
            doc.add_heading("Section 2: Data Table", level=1)
            table = doc.add_table(rows=3, cols=3)
            table.cell(0, 0).text = "Name"
            table.cell(0, 1).text = "Value"
            table.cell(0, 2).text = "Unit"
            table.cell(1, 0).text = "Precision"
            table.cell(1, 1).text = "0.95"
            table.cell(1, 2).text = "ratio"
            table.cell(2, 0).text = "Recall"
            table.cell(2, 1).text = "0.88"
            table.cell(2, 2).text = "ratio"
        elif variant == 2:  # tables
            doc.add_heading("Tables Fixture", 0)
            doc.add_paragraph("Parser benchmark results:")
            table = doc.add_table(rows=5, cols=4)
            headers = ["Document", "Format", "Blocks", "Score"]
            for c, h in enumerate(headers):
                table.cell(0, c).text = h
            rows = [
                ("doc-001", "PDF", "12", "0.91"),
                ("doc-002", "DOCX", "8", "0.87"),
                ("doc-003", "HTML", "6", "0.95"),
                ("doc-004", "TXT", "3", "0.78"),
            ]
            for r, row in enumerate(rows, 1):
                for c, val in enumerate(row):
                    table.cell(r, c).text = val
        elif variant == 3:  # unicode
            doc.add_heading("Unicode Fixture", 0)
            doc.add_paragraph("Arabic: مرحبا بالعالم")
            doc.add_paragraph("Chinese: 中文文本示例")
            doc.add_paragraph("Greek: αβγδεζηθ")
            doc.add_paragraph("Russian: Привет мир")
        elif variant == 4:  # numbered-lists
            doc.add_heading("Lists Fixture", 0)
            doc.add_paragraph("Ordered items:", style="Normal")
            for i, item in enumerate(["First item", "Second item", "Third item"], 1):
                doc.add_paragraph(f"{i}. {item}", style="Normal")
            doc.add_paragraph("Unordered items:", style="Normal")
            for item in ["Alpha", "Beta", "Gamma", "Delta"]:
                doc.add_paragraph(f"- {item}", style="Normal")
        elif variant == 5:  # empty-sections
            doc.add_heading("Empty Sections Fixture", 0)
            for i in range(1, 6):
                doc.add_heading(f"Section {i}", level=1)
                doc.add_paragraph(f"Placeholder content for section {i}.")
        elif variant == 6:  # long-text
            doc.add_heading("Long Text Fixture", 0)
            para = (
                "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
                "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
                "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris. "
            )
            for _ in range(8):
                doc.add_paragraph(para)
        elif variant == 7:  # metadata
            doc.core_properties.author = "AksharaMD Benchmark"
            doc.core_properties.title = "Metadata Fixture"
            doc.core_properties.subject = "Parser coverage testing"
            doc.core_properties.keywords = "aksharamd benchmark synthetic"
            doc.add_heading("Metadata Fixture", 0)
            doc.add_paragraph("This document has core properties (author, title, subject, keywords) set.")
            doc.add_paragraph("Metadata extraction is tested alongside content extraction.")
        elif variant == 8:  # nested-headings
            doc.add_heading("Nested Headings Fixture", 0)
            for i in range(1, 4):
                doc.add_heading(f"Level {i} Heading", level=i)
                doc.add_paragraph(f"Content under level {i} heading.")
                if i < 3:
                    doc.add_heading(f"Level {i} Secondary Heading", level=i)
                    doc.add_paragraph(f"Secondary content under level {i}.")
        elif variant == 9:  # mixed-content
            doc.add_heading("Mixed Content Fixture", 0)
            doc.add_paragraph("This document mixes headings, paragraphs, and a table.")
            doc.add_heading("Overview", level=1)
            doc.add_paragraph("Introductory paragraph with some text.")
            table = doc.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Key"
            table.cell(0, 1).text = "Value"
            table.cell(1, 0).text = "Result"
            table.cell(1, 1).text = "42"
            doc.add_heading("Conclusion", level=1)
            doc.add_paragraph("Final paragraph after the table.")
        else:  # variant == 10: minimal
            doc.add_heading("Minimal", 0)
            doc.add_paragraph("Minimal fixture.")
        doc.save(str(dest))
    except ImportError:
        dest.write_text(
            f"[DOCX variant {variant} skipped: python-docx not installed. "
            "pip install python-docx to generate this fixture.]\n"
        )


def _build_xlsx(dest: Path, variant: int = 1) -> None:
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Benchmark"
        if variant == 1:  # single-sheet
            ws.append(["Document", "Format", "Pages", "Blocks", "Score"])
            for doc, fmt, pages, blocks, score in [
                ("sample-001", "PDF", 4, 12, 0.91),
                ("sample-002", "DOCX", 2, 8, 0.87),
                ("sample-003", "HTML", 1, 6, 0.95),
                ("sample-004", "TXT", 1, 3, 0.78),
            ]:
                ws.append([doc, fmt, pages, blocks, score])
            ws2 = wb.create_sheet("Summary")
            ws2.append(["Metric", "Value"])
            ws2.append(["Total Files", 4])
            ws2.append(["Mean Score", 0.8775])
        elif variant == 2:  # multi-sheet
            ws.title = "Summary"
            ws.append(["Metric", "Value"])
            ws.append(["Total", 10])
            ws.append(["Success", 9])
            ws2 = wb.create_sheet("PDF Results")
            ws2.append(["ID", "Pages", "Blocks", "Elapsed"])
            ws2.append(["pdf-001", 1, 3, 0.12])
            ws3 = wb.create_sheet("Synthetic Results")
            ws3.append(["ID", "Format", "Blocks"])
            ws3.append(["syn-001", "docx", 8])
        elif variant == 3:  # empty-cells
            ws.append(["A", "B", "C", "D"])
            ws.append(["x", None, "z", None])
            ws.append([None, "y", None, "w"])
            ws.append(["a", "b", None, None])
        elif variant == 4:  # numeric-data
            ws.append(["index", "value", "squared", "sqrt"])
            import math
            for i in range(1, 11):
                ws.append([i, i * 1.5, i ** 2, round(math.sqrt(i), 4)])
        elif variant == 5:  # large-table
            ws.append(["ID", "Name", "Category", "Score", "Rank"])
            for i in range(1, 26):
                ws.append([i, f"item-{i:03d}", f"cat-{(i % 3) + 1}", round(0.5 + i * 0.02, 3), i])
        elif variant == 6:  # string-data
            ws.append(["Country", "Capital", "Language", "Region"])
            for row in [
                ("France", "Paris", "French", "Europe"),
                ("Germany", "Berlin", "German", "Europe"),
                ("Japan", "Tokyo", "Japanese", "Asia"),
                ("Brazil", "Brasília", "Portuguese", "South America"),
                ("Egypt", "Cairo", "Arabic", "Africa"),
            ]:
                ws.append(list(row))
        elif variant == 7:  # date-values
            ws.append(["Date", "Event", "Count", "Notes"])
            ws.append(["2024-01-01", "Start", 0, "Baseline"])
            ws.append(["2024-03-15", "Checkpoint", 50, "Midpoint"])
            ws.append(["2024-06-30", "Review", 80, "On track"])
            ws.append(["2024-12-31", "End", 100, "Complete"])
        elif variant == 8:  # unicode
            ws.append(["名前", "値", "単位", "スコア"])
            ws.append(["精度", "0.95", "比率", "A"])
            ws.append(["再現率", "0.88", "比率", "B"])
            ws.append(["アラビア語", "عربي", "语言", "C"])
        elif variant == 9:  # wide-headers
            ws.append(["A", "B", "C", "D", "E", "F", "G", "H", "I", "J"])
            ws.append([1, 2, 3, 4, 5, 6, 7, 8, 9, 10])
            ws.append([11, 12, 13, 14, 15, 16, 17, 18, 19, 20])
        else:  # variant == 10: minimal
            ws.append(["x", "y"])
            ws.append([1, 2])
        wb.save(str(dest))
    except ImportError:
        dest.write_text(
            f"[XLSX variant {variant} skipped: openpyxl not installed. "
            "pip install openpyxl to generate this fixture.]\n"
        )


def _build_pptx(dest: Path, variant: int = 1) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        if variant == 1:  # title-content
            s = prs.slides.add_slide(prs.slide_layouts[0])
            s.shapes.title.text = "AksharaMD Benchmark Deck"
            s.placeholders[1].text = "Public Reproducible Corpus — Synthetic PPTX Fixture"
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            s2.shapes.title.text = "Supported Formats"
            tf = s2.placeholders[1].text_frame
            tf.text = "PDF (34 real files)"
            tf.add_paragraph().text = "DOCX, XLSX, PPTX (synthetic)"
            tf.add_paragraph().text = "HTML, CSV, JSON, XML, TXT, MD, ZIP (synthetic)"
        elif variant == 2:  # multi-slide-5
            titles = [
                "Introduction", "Corpus Overview", "PDF Robustness",
                "Synthetic Formats", "Conclusions"
            ]
            for title in titles:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = title
                s.placeholders[1].text_frame.text = f"Content for: {title}"
        elif variant == 3:  # with-table
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = "Results Table"
            tbl = s.shapes.add_table(
                4, 3, Inches(1), Inches(2), Inches(8), Inches(2)
            ).table
            headers = ["Format", "Files", "Success Rate"]
            data = [("PDF", "34", "94%"), ("DOCX", "10", "100%"), ("CSV", "10", "100%")]
            for c, h in enumerate(headers):
                tbl.cell(0, c).text = h
            for r, row in enumerate(data, 1):
                for c, val in enumerate(row):
                    tbl.cell(r, c).text = val
        elif variant == 4:  # with-bullets
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = "Metrics Recorded"
            tf = s.placeholders[1].text_frame
            tf.text = "Parser success/failure"
            for item in ["Block count and types", "Output character count",
                         "Estimated token count", "Elapsed time per file"]:
                tf.add_paragraph().text = item
        elif variant == 5:  # unicode
            s = prs.slides.add_slide(prs.slide_layouts[0])
            s.shapes.title.text = "ベンチマークデッキ"
            s.placeholders[1].text = "日本語テキスト 中文"
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            s2.shapes.title.text = "عربي"
            s2.placeholders[1].text_frame.text = "مرحبا بالعالم"
        elif variant == 6:  # mixed-layouts
            s1 = prs.slides.add_slide(prs.slide_layouts[0])
            s1.shapes.title.text = "Title Slide"
            s1.placeholders[1].text = "Subtitle text"
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            s2.shapes.title.text = "Content Slide"
            s2.placeholders[1].text_frame.text = "Main content here"
            s3 = prs.slides.add_slide(prs.slide_layouts[1])
            s3.shapes.title.text = "Third Slide"
            s3.placeholders[1].text_frame.text = "More content"
        elif variant == 7:  # long-text
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = "Long Text Slide"
            tf = s.placeholders[1].text_frame
            tf.text = "Lorem ipsum dolor sit amet, consectetur adipiscing elit."
            for sentence in [
                "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.",
                "Ut enim ad minim veniam, quis nostrud exercitation.",
                "Duis aute irure dolor in reprehenderit in voluptate.",
                "Excepteur sint occaecat cupidatat non proident.",
            ]:
                tf.add_paragraph().text = sentence
        elif variant == 8:  # minimal
            s = prs.slides.add_slide(prs.slide_layouts[0])
            s.shapes.title.text = "Minimal"
            s.placeholders[1].text = "Single slide fixture"
        elif variant == 9:  # title-only
            for i in range(1, 4):
                s = prs.slides.add_slide(prs.slide_layouts[0])
                s.shapes.title.text = f"Section {i}"
                s.placeholders[1].text = f"Section {i} subtitle"
        else:  # variant == 10: content-heavy-8
            for i in range(1, 9):
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = f"Topic {i}: Detail"
                tf = s.placeholders[1].text_frame
                tf.text = f"Primary point for topic {i}."
                tf.add_paragraph().text = f"Supporting detail A for topic {i}."
                tf.add_paragraph().text = f"Supporting detail B for topic {i}."
        prs.save(str(dest))
    except ImportError:
        dest.write_text(
            f"[PPTX variant {variant} skipped: python-pptx not installed. "
            "pip install python-pptx to generate this fixture.]\n"
        )


def _build_html(dest: Path, variant: int = 1) -> None:
    if variant == 1:  # headings-paragraphs
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>AksharaMD Benchmark Page</title></head>
<body>
  <h1>AksharaMD Benchmark: HTML Fixture</h1>
  <p>This page is a synthetic HTML fixture for the AksharaMD public benchmark suite.</p>
  <h2>Feature Table</h2>
  <table>
    <thead><tr><th>Feature</th><th>Status</th></tr></thead>
    <tbody>
      <tr><td>Table extraction</td><td>Supported</td></tr>
      <tr><td>Heading detection</td><td>Supported</td></tr>
      <tr><td>Link preservation</td><td>Supported</td></tr>
      <tr><td>Image alt text</td><td>Supported</td></tr>
    </tbody>
  </table>
  <h2>Sample Paragraph</h2>
  <p>The quick brown fox jumps over the lazy dog.
     Pack my box with five dozen liquor jugs.</p>
  <ul><li>Item one</li><li>Item two</li><li>Item three</li></ul>
</body>
</html>
"""
    elif variant == 2:  # tables
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Tables Fixture</title></head>
<body>
  <h1>Tables Fixture</h1>
  <table>
    <thead><tr><th>ID</th><th>Format</th><th>Pages</th><th>Blocks</th><th>Score</th></tr></thead>
    <tbody>
      <tr><td>1</td><td>pdf</td><td>4</td><td>12</td><td>0.91</td></tr>
      <tr><td>2</td><td>docx</td><td>2</td><td>8</td><td>0.87</td></tr>
      <tr><td>3</td><td>html</td><td>1</td><td>6</td><td>0.95</td></tr>
      <tr><td>4</td><td>txt</td><td>1</td><td>3</td><td>0.78</td></tr>
    </tbody>
  </table>
</body>
</html>
"""
    elif variant == 3:  # ordered-unordered-lists
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Lists Fixture</title></head>
<body>
  <h1>Lists Fixture</h1>
  <h2>Ordered List</h2>
  <ol>
    <li>First item</li>
    <li>Second item</li>
    <li>Third item</li>
  </ol>
  <h2>Unordered List</h2>
  <ul>
    <li>Alpha</li>
    <li>Beta</li>
    <li>Gamma</li>
  </ul>
  <h2>Nested List</h2>
  <ul>
    <li>Parent A
      <ul><li>Child A1</li><li>Child A2</li></ul>
    </li>
    <li>Parent B
      <ul><li>Child B1</li></ul>
    </li>
  </ul>
</body>
</html>
"""
    elif variant == 4:  # unicode
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Unicode Fixture</title></head>
<body>
  <h1>Unicode Fixture</h1>
  <p>Arabic: مرحبا بالعالم</p>
  <p>Chinese: 中文文本示例</p>
  <p>Japanese: ベンチマーク</p>
  <p>Russian: Привет мир</p>
  <p>Greek: αβγδε</p>
</body>
</html>
"""
    elif variant == 5:  # nested-divs
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Nested Divs Fixture</title></head>
<body>
  <div class="container">
    <div class="header"><h1>Nested Divs</h1></div>
    <div class="content">
      <div class="section">
        <h2>Section A</h2>
        <div class="subsection"><p>Subsection A1 content.</p></div>
        <div class="subsection"><p>Subsection A2 content.</p></div>
      </div>
      <div class="section">
        <h2>Section B</h2>
        <p>Section B content.</p>
      </div>
    </div>
    <div class="footer"><p>Footer text.</p></div>
  </div>
</body>
</html>
"""
    elif variant == 6:  # metadata-tags
        content = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="description" content="AksharaMD benchmark synthetic HTML fixture">
  <meta name="keywords" content="aksharamd, benchmark, parser, coverage">
  <meta name="author" content="AksharaMD Benchmark">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Metadata Tags Fixture</title>
</head>
<body>
  <h1>Metadata Tags Fixture</h1>
  <p>This page tests metadata tag extraction from HTML documents.</p>
</body>
</html>
"""
    elif variant == 7:  # code-blocks
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Code Blocks Fixture</title></head>
<body>
  <h1>Code Blocks Fixture</h1>
  <h2>Python Example</h2>
  <pre><code>def greet(name: str) -> str:
    return f"Hello, {name}!"

result = greet("World")
print(result)
</code></pre>
  <h2>JSON Example</h2>
  <pre><code>{"key": "value", "count": 42, "enabled": true}
</code></pre>
  <p>Code blocks test parser handling of preformatted content.</p>
</body>
</html>
"""
    elif variant == 8:  # links-images
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Links and Images Fixture</title></head>
<body>
  <h1>Links and Images Fixture</h1>
  <p>Visit <a href="https://github.com/py-pdf/sample-files">py-pdf/sample-files</a> for PDF corpus.</p>
  <p><img src="placeholder.png" alt="Benchmark placeholder image" width="400" height="300"></p>
  <h2>Related Resources</h2>
  <ul>
    <li><a href="https://github.com/run-llama/ParseBench">ParseBench framework</a></li>
    <li><a href="https://huggingface.co/datasets/llamaindex/ParseBench">ParseBench corpus</a></li>
  </ul>
</body>
</html>
"""
    elif variant == 9:  # empty-body
        content = """<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Empty Body Fixture</title></head>
<body></body>
</html>
"""
    else:  # variant == 10: minimal
        content = "<!DOCTYPE html><html><head><title>Minimal</title></head><body><p>Minimal.</p></body></html>\n"

    dest.write_text(content, encoding="utf-8")


def _build_csv(dest: Path, variant: int = 1) -> None:
    if variant == 1:  # simple-tabular
        content = (
            "id,format,pages,blocks,score,notes\n"
            "1,pdf,4,12,0.91,multipage pdflatex\n"
            "2,docx,2,8,0.87,synthetic fixture\n"
            "3,html,1,6,0.95,simple page\n"
            "4,txt,1,3,0.78,plain text\n"
            "5,xlsx,3,15,0.82,spreadsheet with two sheets\n"
            "6,pptx,3,9,0.79,three-slide deck\n"
        )
    elif variant == 2:  # unicode
        content = (
            "name,value,language\n"
            "中文,42,Chinese\n"
            "مرحبا,99,Arabic\n"
            "Привет,17,Russian\n"
            "αβγ,3.14,Greek\n"
        )
    elif variant == 3:  # empty-cells
        content = (
            "a,b,c,d\n"
            "1,,3,\n"
            ",2,,4\n"
            "5,6,,\n"
            ",,7,8\n"
        )
    elif variant == 4:  # large-50rows
        lines = ["id,name,score,category"]
        for i in range(1, 51):
            lines.append(f"{i},item-{i:03d},{0.5 + i * 0.01:.3f},cat-{(i % 5) + 1}")
        content = "\n".join(lines) + "\n"
    elif variant == 5:  # single-column
        content = "value\n1\n2\n3\n4\n5\n6\n7\n8\n9\n10\n"
    elif variant == 6:  # quoted-fields
        content = (
            "name,description,tags\n"
            '"Smith, John","A researcher, author","ai,ml"\n'
            '"Lee, Alice","\"Expert\" in NLP","nlp,parsing"\n'
            '"Brown, Bob","Works on data, models","data,ml"\n'
        )
    elif variant == 7:  # numeric-only
        content = "x,y,z\n"
        for i in range(1, 11):
            content += f"{i},{i * 2.5:.1f},{i ** 2}\n"
    elif variant == 8:  # mixed-types
        content = (
            "id,name,count,ratio,active\n"
            "1,alpha,10,0.95,true\n"
            "2,beta,0,0.0,false\n"
            "3,gamma,100,1.0,true\n"
            "4,delta,,0.5,\n"
        )
    elif variant == 9:  # header-only
        content = "id,name,value,unit,notes\n"
    else:  # variant == 10: minimal
        content = "x,y\n1,2\n"
    dest.write_text(content, encoding="utf-8")


def _build_json(dest: Path, variant: int = 1) -> None:
    if variant == 1:  # flat-object
        data: object = {
            "benchmark": "aksharamd-public-v1",
            "description": "Synthetic JSON fixture for parser coverage testing.",
            "formats_tested": ["pdf", "docx", "xlsx", "pptx", "html", "csv", "json", "xml", "txt", "md", "zip"],
            "results": [
                {"id": f"file-{i:03d}", "format": fmt, "success": True, "blocks": n}
                for i, (fmt, n) in enumerate([
                    ("pdf", 12), ("docx", 8), ("html", 6), ("txt", 3),
                    ("csv", 4), ("json", 2), ("xml", 5), ("md", 7),
                ], start=1)
            ],
            "summary": {"total": 8, "success": 8, "failure": 0, "success_rate": 1.0},
        }
    elif variant == 2:  # nested-object
        data = {
            "config": {
                "version": "2.0",
                "settings": {
                    "timeout": 30,
                    "retries": 3,
                    "output": {"format": "jsonl", "dir": "results/"}
                }
            },
            "corpus": {
                "pdf": {"count": 34, "source": "py-pdf/sample-files"},
                "synthetic": {"count": 100, "formats": 10}
            }
        }
    elif variant == 3:  # array-of-objects
        data = [
            {"id": f"entry-{i:03d}", "format": fmt, "pages": pages, "blocks": blks}
            for i, (fmt, pages, blks) in enumerate([
                ("pdf", 4, 12), ("docx", 2, 8), ("xlsx", 3, 15),
                ("pptx", 3, 9), ("html", 1, 6), ("csv", 1, 4),
                ("json", 1, 2), ("xml", 1, 5), ("txt", 1, 3), ("md", 1, 7),
            ], start=1)
        ]
    elif variant == 4:  # deeply-nested
        data = {
            "level1": {
                "level2": {
                    "level3": {
                        "level4": {
                            "value": "deeply nested value",
                            "count": 42
                        }
                    }
                }
            }
        }
    elif variant == 5:  # unicode
        data = {
            "languages": [
                {"name": "中文", "code": "zh", "sample": "中文文本"},
                {"name": "日本語", "code": "ja", "sample": "日本語テキスト"},
                {"name": "عربي", "code": "ar", "sample": "مرحبا"},
                {"name": "Русский", "code": "ru", "sample": "Привет"},
            ]
        }
    elif variant == 6:  # null-values
        data = {
            "records": [
                {"id": 1, "name": "alpha", "value": 42, "notes": None},
                {"id": 2, "name": None, "value": 0, "notes": "empty"},
                {"id": 3, "name": "gamma", "value": None, "notes": None},
            ]
        }
    elif variant == 7:  # large-array
        data = [
            {"index": i, "label": f"item-{i:03d}", "weight": round(i * 0.05, 3)}
            for i in range(1, 21)
        ]
    elif variant == 8:  # mixed-types
        data = {
            "integer": 42,
            "float": 3.14159,
            "string": "hello world",
            "boolean_true": True,
            "boolean_false": False,
            "null_value": None,
            "array": [1, "two", 3.0, True, None],
            "nested": {"key": "value"}
        }
    elif variant == 9:  # empty-object
        data = {}
    else:  # variant == 10: minimal
        data = {"benchmark": "aksharamd", "version": "2.0"}

    dest.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def _build_xml(dest: Path, variant: int = 1) -> None:
    if variant == 1:  # simple-elements
        content = """<?xml version="1.0" encoding="UTF-8"?>
<benchmark version="1.0">
  <description>Synthetic XML fixture for AksharaMD parser coverage benchmarks.</description>
  <corpus>
    <file id="pdf-001" format="pdf" pages="1">
      <label>minimal-pdflatex</label>
      <expected_outcome>success</expected_outcome>
    </file>
    <file id="syn-001" format="docx" pages="2">
      <label>synthetic-docx</label>
      <expected_outcome>success</expected_outcome>
    </file>
  </corpus>
  <metrics>
    <metric name="success_rate" value="1.0" unit="ratio"/>
    <metric name="mean_blocks" value="7.5" unit="count"/>
  </metrics>
</benchmark>
"""
    elif variant == 2:  # nested-elements
        content = """<?xml version="1.0" encoding="UTF-8"?>
<root>
  <section id="1">
    <title>Section One</title>
    <subsection id="1.1">
      <title>Subsection One</title>
      <paragraph>Nested content in subsection one.</paragraph>
    </subsection>
    <subsection id="1.2">
      <title>Subsection Two</title>
      <paragraph>Nested content in subsection two.</paragraph>
    </subsection>
  </section>
  <section id="2">
    <title>Section Two</title>
    <paragraph>Content in section two.</paragraph>
  </section>
</root>
"""
    elif variant == 3:  # attributes
        content = """<?xml version="1.0" encoding="UTF-8"?>
<catalog xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <item id="001" type="pdf" status="active" priority="high">
    <name>Minimal PDF</name>
    <pages count="1" encrypted="false"/>
  </item>
  <item id="002" type="docx" status="active" priority="medium">
    <name>Synthetic DOCX</name>
    <pages count="2" encrypted="false"/>
  </item>
</catalog>
"""
    elif variant == 4:  # unicode
        content = """<?xml version="1.0" encoding="UTF-8"?>
<languages>
  <language code="zh">中文文本示例</language>
  <language code="ja">日本語テキスト</language>
  <language code="ar">مرحبا بالعالم</language>
  <language code="ru">Привет мир</language>
  <language code="el">αβγδεζ</language>
</languages>
"""
    elif variant == 5:  # mixed-content
        content = """<?xml version="1.0" encoding="UTF-8"?>
<document>
  <intro>This document has <em>mixed</em> content with <strong>inline markup</strong>.</intro>
  <section>
    First some text, then <note>a note</note>, then more text.
  </section>
</document>
"""
    elif variant == 6:  # deep-nesting
        content = """<?xml version="1.0" encoding="UTF-8"?>
<l1>
  <l2>
    <l3>
      <l4>
        <l5>Deeply nested content at level 5.</l5>
      </l4>
    </l3>
  </l2>
</l1>
"""
    elif variant == 7:  # namespaced
        content = """<?xml version="1.0" encoding="UTF-8"?>
<bm:benchmark xmlns:bm="urn:aksharamd:benchmark:v2" xmlns:dc="http://purl.org/dc/elements/1.1/">
  <dc:title>AksharaMD Benchmark</dc:title>
  <dc:creator>AksharaMD</dc:creator>
  <bm:corpus>
    <bm:file bm:id="pdf-001" bm:format="pdf"/>
    <bm:file bm:id="syn-001" bm:format="docx"/>
  </bm:corpus>
</bm:benchmark>
"""
    elif variant == 8:  # cdata
        content = """<?xml version="1.0" encoding="UTF-8"?>
<root>
  <code><![CDATA[
def greet(name: str) -> str:
    return f"Hello, {name}!"
]]></code>
  <json><![CDATA[{"key": "value", "count": 42}]]></json>
  <description>Document with CDATA sections.</description>
</root>
"""
    elif variant == 9:  # empty-elements
        content = """<?xml version="1.0" encoding="UTF-8"?>
<root>
  <empty1/>
  <empty2 attr="value"/>
  <nonempty>content</nonempty>
  <empty3></empty3>
</root>
"""
    else:  # variant == 10: minimal
        content = """<?xml version="1.0" encoding="UTF-8"?>
<root><value>minimal</value></root>
"""
    dest.write_text(content, encoding="utf-8")


def _build_txt(dest: Path, variant: int = 1) -> None:
    if variant == 1:  # paragraphs
        content = (
            "AksharaMD Public Benchmark — Plain Text Fixture\n"
            "================================================\n\n"
            "This file is a plain-text fixture for testing the AksharaMD text parser.\n\n"
            "Section 1: Background\n"
            "---------------------\n"
            "AksharaMD is an open-source document ingestion pipeline that converts 40+\n"
            "file formats into structured Markdown and JSON suitable for RAG workflows.\n\n"
            "Section 2: Supported Formats\n"
            "----------------------------\n"
            "PDF, DOCX, XLSX, PPTX, HTML, EPUB, CSV, JSON, XML, TXT, Markdown,\n"
            "ZIP, TAR, 7z, EML, and more.\n\n"
            "Section 3: Benchmark Scope\n"
            "--------------------------\n"
            "This benchmark measures parser coverage and extraction readiness.\n"
            "It does not measure answer correctness or RAG faithfulness.\n"
        )
    elif variant == 2:  # headings-sections
        content = (
            "CHAPTER 1: INTRODUCTION\n"
            "=======================\n\n"
            "This chapter introduces the benchmark framework.\n\n"
            "1.1 Background\n"
            "--------------\n"
            "Document parsing is a core task in RAG pipelines.\n\n"
            "1.2 Motivation\n"
            "--------------\n"
            "Reproducible benchmarks enable regression detection across releases.\n\n"
            "CHAPTER 2: METHODOLOGY\n"
            "======================\n\n"
            "2.1 Corpus Selection\n"
            "--------------------\n"
            "The corpus covers 10+ file formats from real and synthetic sources.\n\n"
            "2.2 Metrics\n"
            "-----------\n"
            "Parser success rate, block count, and output character count.\n"
        )
    elif variant == 3:  # unicode
        content = (
            "Unicode Text Fixture\n"
            "====================\n\n"
            "Arabic: مرحبا بالعالم\n\n"
            "Chinese: 中文文本示例\n\n"
            "Japanese: ベンチマークデッキ\n\n"
            "Russian: Привет мир\n\n"
            "Greek: αβγδεζηθ\n\n"
            "Emoji text: text followed by unicode symbols • ✓ ✔\n"
        )
    elif variant == 4:  # long-text
        para = (
            "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
            "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
            "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris "
            "nisi ut aliquip ex ea commodo consequat.\n\n"
        )
        content = "Long Text Fixture\n=================\n\n" + para * 12
    elif variant == 5:  # lists
        content = (
            "Lists Fixture\n"
            "=============\n\n"
            "Dashed list:\n"
            "- First item\n"
            "- Second item\n"
            "- Third item\n\n"
            "Asterisk list:\n"
            "* Alpha\n"
            "* Beta\n"
            "* Gamma\n\n"
            "Numbered list:\n"
            "1. One\n"
            "2. Two\n"
            "3. Three\n\n"
            "Indented list:\n"
            "  - Nested A\n"
            "    - Nested A1\n"
            "    - Nested A2\n"
            "  - Nested B\n"
        )
    elif variant == 6:  # minimal
        content = "Minimal plain text fixture.\n"
    elif variant == 7:  # ascii-table
        content = (
            "ASCII Table Fixture\n"
            "===================\n\n"
            "+----------+--------+-------+\n"
            "| Format   | Count  | Score |\n"
            "+----------+--------+-------+\n"
            "| PDF      |     34 |  0.94 |\n"
            "| DOCX     |     10 |  1.00 |\n"
            "| XLSX     |     10 |  1.00 |\n"
            "| HTML     |     10 |  1.00 |\n"
            "| CSV      |     10 |  1.00 |\n"
            "+----------+--------+-------+\n"
        )
    elif variant == 8:  # code-like
        content = (
            "Code-Like Fixture\n"
            "=================\n\n"
            "Python snippet:\n\n"
            "    def compile(path: str) -> Context:\n"
            "        compiler = Compiler()\n"
            "        return compiler.compile(path)\n\n"
            "Shell commands:\n\n"
            "    python benchmarks/build_public_corpus.py\n"
            "    python benchmarks/run_public_benchmark.py --smoke\n\n"
            "JSON fragment:\n\n"
            '    {"version": "2.0", "total": 134}\n'
        )
    elif variant == 9:  # mixed-sections
        content = (
            "Mixed Sections Fixture\n"
            "======================\n\n"
            "PART A: Text\n"
            "------------\n"
            "Standard paragraph with some prose content.\n\n"
            "PART B: Lists\n"
            "-------------\n"
            "- Item one\n"
            "- Item two\n\n"
            "PART C: Data\n"
            "------------\n"
            "x=1, y=2, z=3\n"
            "a=alpha, b=beta\n\n"
            "PART D: Code\n"
            "------------\n"
            "    result = parse(document)\n"
        )
    else:  # variant == 10: whitespace-heavy
        content = (
            "Whitespace Heavy Fixture\n\n\n"
            "First section with extra blank lines.\n\n\n\n"
            "Second section with trailing spaces.   \n\n"
            "Third section with mixed   spacing  inside  words.\n\n\n"
            "End of document.\n"
        )
    dest.write_text(content, encoding="utf-8")


def _build_md(dest: Path, variant: int = 1) -> None:
    if variant == 1:  # full-document
        content = (
            "# AksharaMD Public Benchmark\n\n"
            "Synthetic Markdown fixture for parser coverage testing.\n\n"
            "## Overview\n\n"
            "This corpus tests AksharaMD's ability to extract structured content from\n"
            "a diverse set of document formats.\n\n"
            "## Format Coverage\n\n"
            "| Format | Source | Count |\n"
            "| --- | --- | --- |\n"
            "| PDF | py-pdf/sample-files (CC-BY-SA-4.0) | 34 |\n"
            "| DOCX | synthetic | 10 |\n"
            "| XLSX | synthetic | 10 |\n"
            "| PPTX | synthetic | 10 |\n"
            "| HTML | synthetic | 10 |\n"
            "| CSV | synthetic | 10 |\n"
            "| JSON | synthetic | 10 |\n"
            "| XML | synthetic | 10 |\n"
            "| TXT | synthetic | 10 |\n"
            "| MD | synthetic | 10 |\n"
            "| ZIP | synthetic | 10 |\n\n"
            "## What This Measures\n\n"
            "- Parser success/failure rate\n"
            "- Block structure (headings, paragraphs, tables, code blocks)\n"
            "- Output character count\n"
            "- Estimated token count\n"
        )
    elif variant == 2:  # tables
        content = (
            "# Tables Fixture\n\n"
            "## Simple Table\n\n"
            "| ID | Format | Pages | Blocks |\n"
            "| --- | --- | --- | --- |\n"
            "| 1 | pdf | 4 | 12 |\n"
            "| 2 | docx | 2 | 8 |\n"
            "| 3 | html | 1 | 6 |\n\n"
            "## Alignment Table\n\n"
            "| Left | Center | Right |\n"
            "| :--- | :---: | ---: |\n"
            "| a | b | c |\n"
            "| alpha | beta | gamma |\n"
        )
    elif variant == 3:  # code-blocks
        content = (
            "# Code Blocks Fixture\n\n"
            "## Python\n\n"
            "```python\n"
            "from aksharamd.compiler import Compiler\n\n"
            "ctx = Compiler().compile('document.pdf')\n"
            "print(ctx.document.blocks)\n"
            "```\n\n"
            "## Shell\n\n"
            "```bash\n"
            "python benchmarks/run_public_benchmark.py --smoke\n"
            "```\n\n"
            "## JSON\n\n"
            "```json\n"
            '{"version": "2.0", "corpus_dir": ".public_corpus"}\n'
            "```\n"
        )
    elif variant == 4:  # unicode
        content = (
            "# Unicode Fixture\n\n"
            "## 中文标题\n\n"
            "中文文本示例内容。\n\n"
            "## 日本語タイトル\n\n"
            "日本語のテキストサンプル。\n\n"
            "## عنوان عربي\n\n"
            "محتوى عربي نموذجي.\n\n"
            "## Русский заголовок\n\n"
            "Русский текст.\n"
        )
    elif variant == 5:  # nested-lists
        content = (
            "# Nested Lists Fixture\n\n"
            "- Level 1 Item A\n"
            "  - Level 2 Item A1\n"
            "    - Level 3 Item A1a\n"
            "    - Level 3 Item A1b\n"
            "  - Level 2 Item A2\n"
            "- Level 1 Item B\n"
            "  - Level 2 Item B1\n"
            "  - Level 2 Item B2\n"
            "    - Level 3 Item B2a\n\n"
            "1. Ordered top\n"
            "   1. Ordered sub-item 1\n"
            "   2. Ordered sub-item 2\n"
            "2. Second ordered top\n"
        )
    elif variant == 6:  # headings-h1-h6
        content = (
            "# H1 Heading\n\n"
            "Content under H1.\n\n"
            "## H2 Heading\n\n"
            "Content under H2.\n\n"
            "### H3 Heading\n\n"
            "Content under H3.\n\n"
            "#### H4 Heading\n\n"
            "Content under H4.\n\n"
            "##### H5 Heading\n\n"
            "Content under H5.\n\n"
            "###### H6 Heading\n\n"
            "Content under H6.\n"
        )
    elif variant == 7:  # links-images
        content = (
            "# Links and Images Fixture\n\n"
            "Visit [py-pdf/sample-files](https://github.com/py-pdf/sample-files) for the PDF corpus.\n\n"
            "![Benchmark diagram](placeholder.png)\n\n"
            "## Resources\n\n"
            "- [ParseBench](https://github.com/run-llama/ParseBench)\n"
            "- [AksharaMD](https://github.com/K2alyan/aksharaMD)\n\n"
            "Reference-style [link][ref1] and [another][ref2].\n\n"
            "[ref1]: https://github.com/py-pdf/sample-files\n"
            "[ref2]: https://huggingface.co/datasets/llamaindex/ParseBench\n"
        )
    elif variant == 8:  # blockquotes
        content = (
            "# Blockquotes Fixture\n\n"
            "> This is a blockquote. It tests how the parser handles\n"
            "> quoted content in Markdown.\n\n"
            "> Nested blockquotes:\n"
            "> > Inner level quote.\n"
            "> > > Deepest level.\n\n"
            "Regular paragraph after blockquote.\n\n"
            "> **Bold in quote.** And *italic*.\n"
        )
    elif variant == 9:  # minimal
        content = "# Minimal\n\nMinimal Markdown fixture.\n"
    else:  # variant == 10: long-text
        para = (
            "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
            "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
            "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris.\n\n"
        )
        sections = "".join(
            f"## Section {i}\n\n{para}"
            for i in range(1, 9)
        )
        content = "# Long Text Fixture\n\n" + sections

    dest.write_text(content, encoding="utf-8")


def _build_zip(dest: Path, variant: int = 1) -> None:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        if variant == 1:  # mixed-formats
            zf.writestr("readme.md", "# Mixed Archive\n\nThis ZIP is a synthetic benchmark fixture.\n")
            zf.writestr("data/config.json", '{"version": "1.0", "format": "zip-fixture"}\n')
            zf.writestr("data/notes.txt", "Plain text note inside a ZIP archive.\n")
            zf.writestr(
                "src/hello.py",
                'def greet(name: str) -> str:\n    """Return a greeting."""\n    return f"Hello, {name}!"\n',
            )
            zf.writestr("data/table.csv", "col_a,col_b,col_c\n1,foo,true\n2,bar,false\n3,baz,true\n")
        elif variant == 2:  # single-text
            zf.writestr("document.txt", "Single text file inside a ZIP archive.\nLine two.\nLine three.\n")
        elif variant == 3:  # nested-dirs
            zf.writestr("root.txt", "Root level file.\n")
            zf.writestr("a/file_a.txt", "File in directory a.\n")
            zf.writestr("a/b/file_ab.txt", "File in directory a/b.\n")
            zf.writestr("a/b/c/file_abc.txt", "File in directory a/b/c.\n")
            zf.writestr("x/file_x.txt", "File in directory x.\n")
        elif variant == 4:  # many-small-10
            for i in range(1, 11):
                zf.writestr(f"file-{i:02d}.txt", f"Small file {i} content.\n")
        elif variant == 5:  # unicode-filenames
            zf.writestr("readme.txt", "Archive with international filenames.\n")
            zf.writestr("donnees.txt", "French: données\n")
            zf.writestr("ueberblick.txt", "German: Überblick\n")
            zf.writestr("resume.txt", "Résumé content\n")
        elif variant == 6:  # large-text
            large_content = "Large text file content.\n" * 200
            zf.writestr("large.txt", large_content)
            zf.writestr("index.txt", "Single large text file archive.\n")
        elif variant == 7:  # source-code-tree
            zf.writestr("README.md", "# Project\n\nA simple Python project fixture.\n")
            zf.writestr(
                "src/main.py",
                "#!/usr/bin/env python3\n\"\"\"Main entry point.\"\"\"\n\ndef main() -> None:\n    print('hello')\n",
            )
            zf.writestr("src/utils.py", "def helper(x: int) -> int:\n    return x * 2\n")
            zf.writestr("config/settings.json", '{"debug": false, "timeout": 30}\n')
            zf.writestr("tests/test_main.py", "def test_main():\n    pass\n")
        elif variant == 8:  # config-files
            zf.writestr("config.json", '{"version": "2.0", "enabled": true, "timeout": 30}\n')
            zf.writestr(
                "settings.xml",
                "<?xml version=\"1.0\"?>\n<settings><debug>false</debug><timeout>30</timeout></settings>\n",
            )
            zf.writestr("params.txt", "key=value\nalpha=1\nbeta=2\n")
            zf.writestr("manifest.json", '{"files": ["config.json", "settings.xml", "params.txt"]}\n')
        elif variant == 9:  # empty-mixed
            zf.writestr("nonempty.txt", "This file has content.\n")
            zf.writestr("empty.txt", "")
            zf.writestr("also_nonempty.md", "# Has content\n")
            zf.writestr("also_empty.csv", "")
        else:  # variant == 10: minimal
            zf.writestr("hello.txt", "Hello from ZIP.\n")
    dest.write_bytes(buf.getvalue())


_SYNTHETIC_BUILDERS: dict[str, object] = {
    "docx": _build_docx,
    "xlsx": _build_xlsx,
    "pptx": _build_pptx,
    "html": _build_html,
    "csv": _build_csv,
    "json": _build_json,
    "xml": _build_xml,
    "txt": _build_txt,
    "md": _build_md,
    "zip": _build_zip,
}


def build(
    dry_run: bool = False,
    skip_pdf: bool = False,
    smoke: bool = False,
    max_download_mb: float = 100.0,
) -> dict[str, int]:
    manifest = _load_manifest()
    corpus_root = BENCHMARKS / manifest["corpus_dir"]
    corpus_root.mkdir(parents=True, exist_ok=True)

    pdf_corpus = manifest["pdf_corpus"]
    syn_corpus = manifest["synthetic_corpus"]
    formats = syn_corpus["formats"]
    path_template: str = syn_corpus["local_path_template"]

    variants_per_format = (
        syn_corpus["smoke_variants_per_format"] if smoke else syn_corpus["variants_per_format"]
    )

    smoke_ids: set[str] = set(pdf_corpus["smoke_ids"]) if smoke else set()

    counts: dict[str, int] = {"downloaded": 0, "generated": 0, "skipped": 0, "failed": 0}
    cumulative_bytes = 0

    # ── PDF downloads ──────────────────────────────────────────────────────────
    if not skip_pdf:
        for entry in pdf_corpus["files"]:
            if smoke and entry["id"] not in smoke_ids:
                continue

            dest = corpus_root / entry["local_path"]
            print(f"[{entry['id']}] {entry['label']} (pdf)")

            if not dest.exists() and max_download_mb > 0:
                size_hint = _get_remote_size(entry["url"])
                if size_hint is not None:
                    projected_mb = (cumulative_bytes + size_hint) / 1_000_000
                    if projected_mb > max_download_mb:
                        label = entry.get("note") or entry["label"]
                        print(f"  skip (size guard: {size_hint / 1e6:.1f} MB would exceed "
                              f"{max_download_mb} MB limit): {label}")
                        counts["skipped"] += 1
                        continue

            ok, bytes_dl = _download(entry["url"], dest, dry_run=dry_run)
            cumulative_bytes += bytes_dl
            if ok:
                counts["downloaded"] += 1
            else:
                counts["failed"] += 1
    else:
        print("Skipping PDF downloads (--skip-pdf)")

    # ── Synthetic generation ───────────────────────────────────────────────────
    for fmt in formats:
        builder = _SYNTHETIC_BUILDERS.get(fmt)
        for v in range(1, variants_per_format + 1):
            local_path = path_template.format(format=fmt, variant=v, ext=fmt)
            dest = corpus_root / local_path
            print(f"[syn-{fmt}-{v:02d}] {fmt} variant {v}")

            if dest.exists():
                print(f"  skip (cached): {dest.name}")
                counts["skipped"] += 1
                continue

            if dry_run:
                print(f"  [dry-run] would generate: {dest}")
                counts["generated"] += 1
                continue

            if builder is None:
                print(f"  warning: no builder for format '{fmt}' — skipping")
                counts["skipped"] += 1
                continue

            dest.parent.mkdir(parents=True, exist_ok=True)
            try:
                builder(dest, v)  # type: ignore[operator]
                print(f"  generated ({dest.stat().st_size:,} B): {dest.name}")
                counts["generated"] += 1
            except Exception as exc:
                print(f"  error generating {dest.name}: {exc}")
                counts["failed"] += 1

    return counts


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--dry-run", action="store_true", help="Print actions without writing files")
    parser.add_argument("--skip-pdf", action="store_true", help="Skip downloading PDF files")
    mode = parser.add_mutually_exclusive_group()
    mode.add_argument(
        "--smoke", action="store_true",
        help="10 smoke PDFs + 1 variant per format (20 files total)"
    )
    mode.add_argument(
        "--full", action="store_true",
        help="All 34 PDFs + 10 variants per format (134 files total, default)"
    )
    parser.add_argument(
        "--max-download-mb", type=float, default=100.0, metavar="M",
        help="Skip PDFs that would push cumulative download past M MB (0 = unlimited, default: 100)"
    )
    args = parser.parse_args()

    mode_label = "smoke" if args.smoke else "full"
    print(f"Building corpus [{mode_label} mode] into: {BENCHMARKS / '.public_corpus'}\n")

    counts = build(
        dry_run=args.dry_run,
        skip_pdf=args.skip_pdf,
        smoke=args.smoke,
        max_download_mb=args.max_download_mb,
    )

    print(
        f"\nDone. downloaded={counts['downloaded']}  "
        f"generated={counts['generated']}  "
        f"skipped={counts['skipped']}  "
        f"failed={counts['failed']}"
    )
    if counts["failed"] > 0:
        sys.exit(1)


if __name__ == "__main__":
    main()
