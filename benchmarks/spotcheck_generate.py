#!/usr/bin/env python3
"""
Auto-generate spotcheck_facts.yaml from an existing document corpus.

Default mode uses heuristic extraction (no API key needed).
With --llm, uses Claude for higher-quality facts (needs ANTHROPIC_API_KEY).

Usage:
    python -m benchmarks.spotcheck_generate \
        --corpus C:/Users/kalya/Downloads/benchmark_corpus \
        --out benchmarks/spotcheck_facts.yaml \
        --n 2

    # With LLM-quality facts:
    ANTHROPIC_API_KEY=sk-... python -m benchmarks.spotcheck_generate \
        --corpus C:/Users/kalya/Downloads/benchmark_corpus --llm
"""
from __future__ import annotations

import argparse
import re
import sys
import time
from collections import defaultdict
from pathlib import Path

import yaml

_MODEL = "claude-haiku-4-5-20251001"

DEFAULT_FORMATS = [
    "pdf", "docx", "pptx", "html", "xlsx", "epub",
    "txt", "json", "jsonl", "xml", "csv", "ipynb",
    "rss", "atom", "msg",
]

_MIN_CHARS = 300


# ── naive text extraction ─────────────────────────────────────────────────────

def _naive_text(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    try:
        if ext == "pdf":
            import fitz
            doc = fitz.open(str(path))
            text = "".join(p.get_text() for p in doc)
            doc.close()
            return text

        if ext in ("txt", "html", "htm", "xml", "rss", "atom", "csv",
                   "json", "jsonl", "md"):
            return path.read_text(encoding="utf-8", errors="replace")

        if ext == "docx":
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == "pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            return "\n".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes if hasattr(shape, "text")
            )

        if ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            parts.append(str(cell.value))
            return "\n".join(parts)

        if ext == "epub":
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book = epub.read_epub(str(path), options={"ignore_ncx": True})
            parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                parts.append(BeautifulSoup(item.get_content(), "html.parser").get_text())
            return "\n".join(parts)

        if ext == "ipynb":
            import json as _json
            nb = _json.loads(path.read_bytes())
            parts = []
            for cell in nb.get("cells", []):
                src = "".join(cell.get("source", []))
                if src.strip():
                    parts.append(src)
            return "\n\n".join(parts)

        if ext in ("msg", "eml"):
            return path.read_text(encoding="utf-8", errors="replace")[:5000]

    except Exception:
        pass
    return ""


# ── heuristic fact extraction ─────────────────────────────────────────────────

# Patterns that indicate markup/metadata — skip these entirely
_MARKUP_RE = re.compile(
    r'<[^>]+>'                  # HTML/XML tags
    r'|xmlns[:\s]'              # XML namespace declarations
    r'|encoding="'              # XML encoding declaration
    r'|charset='                # HTML charset
    r'|http[s]?://'            # URLs
    r'|www\.'                   # URLs
    r'|\{[\w:]+\}'              # namespace URIs
    r'|^[\s\W]+$'              # lines with no real words
)

def _is_markup(s: str) -> bool:
    return bool(_MARKUP_RE.search(s)) or s.count('<') + s.count('>') + s.count('/') > 3


def _extract_facts_heuristic(text: str) -> list[dict]:
    """
    Extract specific, verifiable content facts — no markup, no URLs, no raw syntax.
    Works on plain prose after naive extraction.
    """
    facts = []
    seen: set[str] = set()

    def add(snippet: str, keywords: list[str]) -> None:
        snippet = re.sub(r'\s+', ' ', snippet).strip()[:120]
        if len(snippet) < 15 or _is_markup(snippet):
            return
        key = snippet.lower()[:40]
        if key in seen:
            return
        seen.add(key)
        entry: dict = {"text": snippet}
        kws = [k for k in keywords if len(k) > 3 and k.lower() in text.lower()][:3]
        if kws:
            entry["keywords"] = kws
        facts.append(entry)

    # Split into lines, skip blank/markup lines, work on clean prose
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    clean_lines = [ln for ln in lines if not _is_markup(ln) and len(ln) > 20]

    # 1. Lines containing a year (4-digit) or specific number — high verifiability
    for ln in clean_lines[:200]:
        if re.search(r'\b(19|20)\d{2}\b', ln) or re.search(r'\b\d{2,}\b', ln):
            nums = re.findall(r'\b[\d,\.]+\b', ln)
            words = [w for w in ln.split() if w.isalpha() and len(w) > 4]
            add(ln[:100], nums[:2] + words[:1])
            if len(facts) >= 3:
                break

    # 2. Multi-word proper nouns (Title Case sequences, 2–5 words)
    proper_re = re.compile(r'(?:[A-Z][a-z]{2,}\s){1,4}[A-Z][a-z]{2,}')
    for m in proper_re.finditer('\n'.join(clean_lines[:100])):
        noun = m.group().strip()
        if not _is_markup(noun):
            add(noun, [w for w in noun.split() if len(w) > 4])
        if len(facts) >= 5:
            break

    # 3. Short, clean heading-style lines (5–12 words, no punctuation clutter)
    for ln in clean_lines[:80]:
        words = ln.split()
        if 4 <= len(words) <= 12 and ln[0].isupper() and not ln.endswith(','):
            add(ln, [w for w in words if w.isalpha() and len(w) > 5][:2])
        if len(facts) >= 5:
            break

    return facts[:5]


# ── LLM fact extraction ───────────────────────────────────────────────────────

def _extract_facts_llm(text: str, filename: str) -> list[dict]:
    import anthropic
    client = anthropic.Anthropic()

    excerpt = text[:4000]
    prompt = f"""You are reviewing a document to create a fact-checking test.

Document: {filename}
Excerpt:
---
{excerpt}
---

Identify 3 to 5 specific, verifiable facts. Each fact should contain a specific number, name, or phrase that can be checked by string search. No vague summaries.

Respond in YAML only:
facts:
  - text: "exact quote or specific claim"
    keywords: ["word1", "word2"]
  - text: "another fact"
    keywords: ["word1"]
"""
    msg = client.messages.create(
        model=_MODEL, max_tokens=512,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = msg.content[0].text.strip()
    if raw.startswith("```"):
        raw = "\n".join(raw.split("\n")[1:]).rstrip("`").strip()
    try:
        return yaml.safe_load(raw).get("facts", [])
    except Exception:
        return []


# ── main ─────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Auto-generate spotcheck facts YAML from corpus")
    parser.add_argument("--corpus", required=True,
                        help="Path to benchmark corpus directory")
    parser.add_argument("--out", default="benchmarks/spotcheck_facts.yaml",
                        help="Output YAML path")
    parser.add_argument("--n", type=int, default=2,
                        help="Files per format (default: 2)")
    parser.add_argument("--formats", nargs="+", default=DEFAULT_FORMATS)
    parser.add_argument("--llm", action="store_true",
                        help="Use Claude for fact extraction (needs ANTHROPIC_API_KEY)")
    args = parser.parse_args()

    if args.llm:
        try:
            import anthropic
            anthropic.Anthropic().models.list()
        except Exception as e:
            print(f"ERROR: LLM not available: {e}", file=sys.stderr)
            sys.exit(1)

    corpus = Path(args.corpus)
    if not corpus.exists():
        print(f"ERROR: corpus not found: {corpus}", file=sys.stderr)
        sys.exit(1)

    by_ext: dict[str, list[Path]] = defaultdict(list)
    for p in corpus.rglob("*"):
        if p.is_file() and p.suffix:
            ext = p.suffix.lower().lstrip(".")
            if ext in args.formats:
                by_ext[ext].append(p)

    documents = []
    total = sum(min(len(v), args.n) for v in by_ext.values())
    done = 0

    for fmt in args.formats:
        for path in by_ext.get(fmt, [])[:args.n]:
            done += 1
            print(f"[{done}/{total}] {fmt:<6} {path.name[:50]} … ", end="", flush=True)

            text = _naive_text(path)
            if len(text.strip()) < _MIN_CHARS:
                print("skip (too short)")
                continue

            if args.llm:
                facts = _extract_facts_llm(text, path.name)
                time.sleep(0.3)
            else:
                facts = _extract_facts_heuristic(text)

            if not facts:
                print("skip (no facts extracted)")
                continue

            documents.append({
                "path": str(path),
                "description": f"{fmt.upper()} — {path.name}",
                "facts": facts,
            })
            print(f"{len(facts)} facts")

    # Sanitise: remove non-BMP characters that confuse PyYAML on Windows
    def _clean(obj):
        if isinstance(obj, str):
            return obj.encode("utf-8", "replace").decode("utf-8").replace("�", "?")
        if isinstance(obj, dict):
            return {k: _clean(v) for k, v in obj.items()}
        if isinstance(obj, list):
            return [_clean(i) for i in obj]
        return obj

    documents = _clean(documents)

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "w", encoding="utf-8") as f:
        mode = "llm" if args.llm else "heuristic"
        f.write(f"# Auto-generated by spotcheck_generate.py  (mode: {mode})\n")
        f.write(f"# Corpus: {args.corpus}  |  Files per format: {args.n}\n\n")
        yaml.dump({"documents": documents}, f, allow_unicode=True,
                  default_flow_style=False, sort_keys=False)

    print(f"\nWrote {len(documents)} documents -> {out}")
    print(f"Run:  python -m benchmarks.spotcheck --facts {out}")


if __name__ == "__main__":
    main()
