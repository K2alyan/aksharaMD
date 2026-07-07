"""
Golden-output tests: compile a minimal, known-content fixture through the full
pipeline and assert stable structural invariants.

Design choices:
  - Fixtures are created programmatically (no committed binaries) so tests stay
    self-contained and CI-reproducible without large file downloads.
  - Assertions check structural invariants (block types, key content, no errors)
    rather than exact string equality, which would be brittle across whitespace
    normalisation or minor optimiser changes.
  - Each test records block counts so accidental truncation is detected even
    when content assertions pass.

To add a new format: add a fixture factory function and a test_<fmt>_golden() test.
"""
from __future__ import annotations

import json
import textwrap
from pathlib import Path

import pytest

from aksharamd.compiler import Compiler
from aksharamd.models.block import BlockType

# ── helpers ───────────────────────────────────────────────────────────────────

def _compile(path: Path, tmp_path: Path) -> tuple[str, object]:
    out = str(tmp_path / "out")
    c = Compiler(output_dir=out)
    return c.compile_to_string(str(path))


def _block_types(doc) -> list[str]:
    if doc is None:
        return []
    return [b.type.value for b in doc.blocks]


def _all_text(doc) -> str:
    if doc is None:
        return ""
    return " ".join(b.content for b in doc.blocks)


# ── HTML ─────────────────────────────────────────────────────────────────────

def test_html_golden(tmp_path):
    f = tmp_path / "doc.html"
    f.write_text(textwrap.dedent("""\
        <!DOCTYPE html><html><head><title>Golden HTML</title></head>
        <body>
          <h1>Section One</h1>
          <p>AksharaMD golden test paragraph with known content.</p>
          <h2>Section Two</h2>
          <ul><li>Alpha</li><li>Beta</li><li>Gamma</li></ul>
          <table>
            <thead><tr><th>Name</th><th>Score</th></tr></thead>
            <tbody>
              <tr><td>Alice</td><td>95</td></tr>
              <tr><td>Bob</td><td>87</td></tr>
            </tbody>
          </table>
        </body></html>
    """), encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    types = _block_types(doc)
    assert BlockType.HEADING.value in types
    assert BlockType.PARAGRAPH.value in types
    assert BlockType.LIST.value in types
    assert BlockType.TABLE.value in types

    content = _all_text(doc)
    assert "Section One" in content
    assert "AksharaMD golden test paragraph" in content
    assert "Alpha" in content
    assert "Alice" in content
    assert "95" in content

    assert len(doc.blocks) >= 5, f"Expected >=5 blocks, got {len(doc.blocks)}"


# ── Markdown ──────────────────────────────────────────────────────────────────

def test_markdown_golden(tmp_path):
    f = tmp_path / "doc.md"
    f.write_text(textwrap.dedent("""\
        # Golden Markdown

        Paragraph with **bold** and *italic* text.

        ## Second Heading

        - item one
        - item two
        - item three

        | Col A | Col B |
        |-------|-------|
        | X     | 1     |
        | Y     | 2     |
    """), encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    types = _block_types(doc)
    assert BlockType.HEADING.value in types

    content = _all_text(doc)
    assert "Golden Markdown" in content
    assert "item one" in content

    assert len(doc.blocks) >= 3, f"Expected >=3 blocks, got {len(doc.blocks)}"


# ── CSV ───────────────────────────────────────────────────────────────────────

def test_csv_golden(tmp_path):
    f = tmp_path / "data.csv"
    rows = [
        ["product", "quantity", "price"],
        ["widget", "100", "9.99"],
        ["gadget", "50", "24.95"],
        ["doohickey", "200", "4.50"],
    ]
    f.write_text("\n".join(",".join(r) for r in rows) + "\n", encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "product" in content.lower() or "widget" in content
    assert "9.99" in content or "24.95" in content

    assert len(doc.blocks) >= 1


# ── JSON ──────────────────────────────────────────────────────────────────────

def test_json_golden(tmp_path):
    f = tmp_path / "record.json"
    payload = {
        "title": "Golden JSON Document",
        "author": "AksharaMD Test Suite",
        "version": 42,
        "tags": ["golden", "test", "structured"],
        "nested": {"key": "value", "count": 3},
    }
    f.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "Golden JSON Document" in content
    assert "AksharaMD Test Suite" in content

    assert len(doc.blocks) >= 2


# ── XML ───────────────────────────────────────────────────────────────────────

def test_xml_golden(tmp_path):
    f = tmp_path / "config.xml"
    f.write_text(textwrap.dedent("""\
        <?xml version="1.0" encoding="UTF-8"?>
        <configuration>
          <title>Golden XML</title>
          <database>
            <host>localhost</host>
            <port>5432</port>
            <name>aksharamd_test</name>
          </database>
          <feature enabled="true">golden-output</feature>
        </configuration>
    """), encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "Golden XML" in content or "golden" in content.lower()
    assert "localhost" in content or "5432" in content

    assert len(doc.blocks) >= 2


# ── Plain text ────────────────────────────────────────────────────────────────

def test_text_golden(tmp_path):
    f = tmp_path / "notes.txt"
    f.write_text(textwrap.dedent("""\
        Meeting Notes — Q3 Review

        Attendees: Alice, Bob, Carol

        Action items:
        1. Alice to update the roadmap by Friday.
        2. Bob to review the security audit.
        3. Carol to draft the release announcement.

        Next meeting: same time next week.
    """), encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "Alice" in content
    assert "security audit" in content

    assert len(doc.blocks) >= 1


# ── EML ───────────────────────────────────────────────────────────────────────

def test_eml_golden(tmp_path):
    f = tmp_path / "message.eml"
    f.write_text(textwrap.dedent("""\
        From: alice@example.com
        To: bob@example.com
        Subject: Golden EML Test
        Date: Mon, 07 Jul 2026 10:00:00 +0000
        MIME-Version: 1.0
        Content-Type: text/plain; charset=utf-8

        This is the golden EML body text.
        It contains multiple lines for the parser to handle.

        Best,
        Alice
    """), encoding="utf-8")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "golden EML body text" in content
    assert "Alice" in content or "alice@example.com" in content

    assert len(doc.blocks) >= 2


# ── ZIP (listing only) ────────────────────────────────────────────────────────

def test_zip_golden(tmp_path):
    import zipfile

    f = tmp_path / "archive.zip"
    with zipfile.ZipFile(str(f), "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("readme.txt", "Golden ZIP readme content.\nLine two.")
        zf.writestr("data.json", '{"key": "golden_value"}')
        zf.writestr("subdir/notes.md", "# Subdir Note\n\nNested content.")

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "readme.txt" in content
    assert "data.json" in content

    assert len(doc.blocks) >= 2


# ── DOCX ──────────────────────────────────────────────────────────────────────

def test_docx_golden(tmp_path):
    pytest.importorskip("docx", reason="python-docx not installed")
    from docx import Document as DocxDocument

    f = tmp_path / "report.docx"
    d = DocxDocument()
    d.add_heading("Golden DOCX Report", level=1)
    d.add_paragraph("This is the introduction paragraph for the golden test.")
    d.add_heading("Section 2: Data", level=2)
    d.add_paragraph("The data section contains structured information.")
    table = d.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "100"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "200"
    d.save(str(f))

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    types = _block_types(doc)
    assert BlockType.HEADING.value in types

    content = _all_text(doc)
    assert "Golden DOCX Report" in content
    assert "introduction paragraph" in content
    assert "Alpha" in content

    assert len(doc.blocks) >= 4


# ── XLSX ──────────────────────────────────────────────────────────────────────

def test_xlsx_golden(tmp_path):
    pytest.importorskip("openpyxl", reason="openpyxl not installed")
    import openpyxl

    f = tmp_path / "budget.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Q3 Budget"
    ws.append(["Category", "Budget", "Actual", "Variance"])
    ws.append(["Engineering", 100000, 95000, 5000])
    ws.append(["Marketing", 50000, 52000, -2000])
    ws.append(["Operations", 75000, 74000, 1000])
    ws2 = wb.create_sheet("Notes")
    ws2["A1"] = "Golden XLSX test fixture"
    wb.save(str(f))

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "Q3 Budget" in content or "Budget" in content
    assert "Engineering" in content or "100000" in content

    assert len(doc.blocks) >= 2


# ── PPTX ─────────────────────────────────────────────────────────────────────

def test_pptx_golden(tmp_path):
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation

    f = tmp_path / "deck.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # title + content

    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Golden PPTX Deck"
    slide1.placeholders[1].text = "Introduction slide for the golden test."

    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Key Findings"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Finding one: extraction is stable"
    tf.add_paragraph().text = "Finding two: structure is preserved"

    prs.save(str(f))

    text, ctx = _compile(f, tmp_path)
    assert not ctx.validation.errors, ctx.validation.errors
    doc = ctx.document

    content = _all_text(doc)
    assert "Golden PPTX" in content
    assert "Finding" in content or "introduction" in content.lower()

    assert len(doc.blocks) >= 2


# ── Safe-mode: subprocess and ML parsers are blocked ─────────────────────────

def test_safe_mode_blocks_audio(tmp_path):
    """In safe mode, audio files must fail with SAFE_MODE_BLOCKED, not attempt ML inference."""
    f = tmp_path / "clip.mp3"
    f.write_bytes(b"\xff\xfb" + b"\x00" * 128)  # minimal fake MP3 header

    out = str(tmp_path / "out")
    c = Compiler(output_dir=out, safe_mode=True)
    _, ctx = c.compile_to_string(str(f))

    error_codes = [e.code for e in ctx.validation.errors]
    assert "SAFE_MODE_BLOCKED" in error_codes


def test_safe_mode_blocks_legacy_office(tmp_path):
    """In safe mode, .doc files must fail with SAFE_MODE_BLOCKED, not invoke soffice."""
    f = tmp_path / "old.doc"
    f.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 500)  # OLE magic header

    out = str(tmp_path / "out")
    c = Compiler(output_dir=out, safe_mode=True)
    _, ctx = c.compile_to_string(str(f))

    error_codes = [e.code for e in ctx.validation.errors]
    assert "SAFE_MODE_BLOCKED" in error_codes


def test_safe_mode_blocks_url(tmp_path):
    """In safe mode, http:// sources must be rejected before any network call."""
    out = str(tmp_path / "out")
    c = Compiler(output_dir=out, safe_mode=True)
    _, ctx = c.compile_to_string("https://example.com/doc.pdf")

    error_codes = [e.code for e in ctx.validation.errors]
    assert "SAFE_MODE_BLOCKED" in error_codes


def test_safe_mode_allows_local_text(tmp_path):
    """In safe mode, plain local text files must still parse successfully."""
    f = tmp_path / "safe.txt"
    f.write_text("Safe mode test content.\nSecond line.", encoding="utf-8")

    out = str(tmp_path / "out")
    c = Compiler(output_dir=out, safe_mode=True)
    text, ctx = c.compile_to_string(str(f))

    assert not ctx.validation.errors, ctx.validation.errors
    assert "Safe mode test content" in text
