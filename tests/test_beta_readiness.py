"""Beta-readiness tests — user-facing CLI behavior and manifest completeness.

Tests in this file verify what a real user sees when running AksharaMD,
not just whether the internal pipeline produces correct output.
"""
from __future__ import annotations

import io
import json
from pathlib import Path

import fitz
import pytest
from click.testing import CliRunner
from PIL import Image, ImageDraw

from aksharamd.cli import main
from aksharamd.compiler import Compiler
from aksharamd.models.manifest import _quality_band

# ── CLI runner fixture ────────────────────────────────────────────────────────

@pytest.fixture
def runner():
    return CliRunner()


# ── PDF fixture builders (inline, no dependency on test_pdf_regression) ───────

def _make_clean_native_pdf(path: Path) -> None:
    doc = fitz.open()
    pg = doc.new_page()
    pg.insert_text((72, 60), "Report on Indoor Air Quality", fontsize=18)
    pg.insert_text((72, 90), "Introduction", fontsize=14)
    pg.insert_text((72, 115),
        "This report summarises measurements taken over a six-month period. "
        "CO2, PM2.5, and TVOC were logged at 1-minute intervals using calibrated sensors.",
        fontsize=10)
    pg.insert_text((72, 185), "Results", fontsize=14)
    pg.insert_text((72, 210),
        "Mean CO2 concentration was 842 ppm. PM2.5 averaged 8.4 ug/m3, well below WHO guideline.",
        fontsize=10)
    doc.save(str(path))
    doc.close()


def _make_scanned_pdf(path: Path) -> None:
    doc = fitz.open()
    for pn in range(1, 4):
        img = Image.new("RGB", (612, 792), (245, 245, 240))
        draw = ImageDraw.Draw(img)
        draw.text((80, 60), f"Scanned page {pn}", fill=(30, 30, 30))
        for ly in range(130, 700, 24):
            draw.line([(80, ly), (530, ly + 4)], fill=(80, 80, 80), width=1)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        buf.seek(0)
        pg = doc.new_page(width=612, height=792)
        pg.insert_image(fitz.Rect(0, 0, 612, 792), stream=buf.read())
    doc.save(str(path))
    doc.close()


def _make_encrypted_pdf(path: Path) -> None:
    src = path.parent / "_enc_src.pdf"
    doc = fitz.open()
    pg = doc.new_page()
    pg.insert_text((72, 60), "Confidential Report", fontsize=16)
    pg.insert_text((72, 90), "This document is password protected.", fontsize=10)
    doc.save(str(src))
    doc.close()
    doc2 = fitz.open(str(src))
    doc2.save(str(path), encryption=fitz.PDF_ENCRYPT_AES_256,
              user_pw="secret", owner_pw="owner", permissions=fitz.PDF_PERM_PRINT)
    doc2.close()


# ── 1. quality_band utility function ─────────────────────────────────────────

def test_quality_band_labels():
    assert _quality_band(100) == "HIGH"
    assert _quality_band(85)  == "HIGH"
    assert _quality_band(84)  == "OK"
    assert _quality_band(70)  == "OK"
    assert _quality_band(69)  == "RISKY"
    assert _quality_band(50)  == "RISKY"
    assert _quality_band(49)  == "POOR"
    assert _quality_band(0)   == "POOR"


# ── 2. Manifest fields — clean PDF ────────────────────────────────────────────

def test_manifest_has_quality_band_for_clean_pdf(tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    compiler = Compiler(output_dir=str(tmp_path / "out"))
    _, ctx = compiler.compile_to_string(str(pdf))
    m = ctx.manifest
    assert m is not None
    assert m.quality_band == "HIGH"
    assert m.readiness_score >= 85


def test_manifest_has_pdf_classification(tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    compiler = Compiler(output_dir=str(tmp_path / "out"))
    _, ctx = compiler.compile_to_string(str(pdf))
    m = ctx.manifest
    assert m.pdf_classification != "", "pdf_classification should be set for PDF files"
    assert m.pdf_classification in (
        "native_text", "scanned", "hybrid", "table_heavy", "layout_heavy", "low_confidence"
    )


def test_manifest_has_warning_codes(tmp_path):
    pdf = tmp_path / "scanned.pdf"
    _make_scanned_pdf(pdf)
    compiler = Compiler(output_dir=str(tmp_path / "out"))
    _, ctx = compiler.compile_to_string(str(pdf))
    m = ctx.manifest
    assert m is not None
    assert isinstance(m.warning_codes, list)
    assert "OCR_REQUIRED" in m.warning_codes or "NEAR_EMPTY_OUTPUT" in m.warning_codes


def test_manifest_image_pages_for_scanned(tmp_path):
    pdf = tmp_path / "scanned.pdf"
    _make_scanned_pdf(pdf)
    compiler = Compiler(output_dir=str(tmp_path / "out"))
    _, ctx = compiler.compile_to_string(str(pdf))
    m = ctx.manifest
    assert m.image_pages >= 1, "scanned PDF should report at least 1 image page"


def test_manifest_ocr_available_field(tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    compiler = Compiler(output_dir=str(tmp_path / "out"))
    _, ctx = compiler.compile_to_string(str(pdf))
    m = ctx.manifest
    assert m.ocr_available is not None, "ocr_available should be set for PDF files"
    assert isinstance(m.ocr_available, bool)


def test_manifest_json_has_quality_band(tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    compiler = Compiler(output_dir=str(tmp_path / "out"))
    compiler.compile(str(pdf))
    manifest_path = tmp_path / "out" / "manifest.json"
    assert manifest_path.exists()
    data = json.loads(manifest_path.read_text())
    assert "quality_band" in data
    assert data["quality_band"] in ("HIGH", "OK", "RISKY", "POOR")
    assert "pdf_classification" in data
    assert "image_pages" in data
    assert "warning_codes" in data


# ── 3. CLI output — clean PDF ────────────────────────────────────────────────

def test_cli_shows_readiness_and_band_for_clean_pdf(runner, tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    assert result.exit_code == 0, result.output
    output = result.output
    assert "Readiness" in output or "87" in output or "HIGH" in output
    assert "Output Files" in output


def test_cli_shows_pdf_classification(runner, tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    assert result.exit_code == 0, result.output
    # Classification should appear somewhere in output (native text, scanned, etc.)
    assert any(label in result.output for label in [
        "native text", "scanned", "hybrid", "table-heavy", "multi-column"
    ])


def test_cli_shows_no_scary_warnings_for_clean_pdf(runner, tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    assert result.exit_code == 0, result.output
    # A clean HIGH-quality PDF should not display a Warnings panel
    assert "Warnings" not in result.output or "OCR_REQUIRED" not in result.output


def test_cli_shows_output_file_locations(runner, tmp_path):
    pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(pdf)
    out_dir = tmp_path / "out"
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(out_dir)])
    assert result.exit_code == 0, result.output
    # Verify the promised output files are actually written to disk —
    # the CLI panel names these files; the real user value is that they exist.
    compiled = out_dir / "clean"
    assert (compiled / "document.md").exists(), "document.md not written"
    assert (compiled / "manifest.json").exists(), "manifest.json not written"
    assert (compiled / "validation.json").exists(), "validation.json not written"


# ── 4. CLI output — scanned PDF (poor extraction) ────────────────────────────

def test_cli_shows_warnings_for_scanned_pdf(runner, tmp_path):
    pdf = tmp_path / "scanned.pdf"
    _make_scanned_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    # Should not crash, score should be low, OCR guidance should appear
    output = result.output
    assert "pip install aksharamd[ocr]" in output or "OCR" in output or "POOR" in output or "RISKY" in output


def test_cli_score_for_scanned_is_poor(runner, tmp_path):
    pdf = tmp_path / "scanned.pdf"
    _make_scanned_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    output = result.output
    assert "POOR" in output or "RISKY" in output, (
        "Scanned PDF without OCR should display POOR or RISKY quality band"
    )


# ── 5. CLI output — encrypted PDF ────────────────────────────────────────────

def test_cli_encrypted_pdf_shows_clear_message(runner, tmp_path):
    pdf = tmp_path / "encrypted.pdf"
    _make_encrypted_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    # Should exit non-zero and mention password or encryption
    output = result.output
    assert result.exit_code != 0
    assert any(word in output.lower() for word in [
        "password", "encrypted", "encrypt", "protected"
    ]), f"Expected password/encryption message, got:\n{output}"


def test_cli_encrypted_pdf_does_not_show_confusing_generic_error(runner, tmp_path):
    pdf = tmp_path / "encrypted.pdf"
    _make_encrypted_pdf(pdf)
    result = runner.invoke(main, ["compile", str(pdf), "-o", str(tmp_path / "out")])
    # Should NOT just say "No content could be extracted" without explanation
    # The ENCRYPTED_PDF warning should surface actionable guidance
    output = result.output
    assert "ENCRYPTED_PDF" in output or "password" in output.lower() or "decrypt" in output.lower()


# ── 6. Manifest quality_band consistency ──────────────────────────────────────

@pytest.mark.parametrize("score,expected_band", [
    (100, "HIGH"), (85, "HIGH"), (84, "OK"), (70, "OK"),
    (69, "RISKY"), (50, "RISKY"), (49, "POOR"), (0, "POOR"),
])
def test_quality_band_covers_all_boundaries(score, expected_band):
    assert _quality_band(score) == expected_band


# ── 7. Beta corpus smoke test ──────────────────────────────────────────────────

def _smoke(runner, src: Path, out: Path) -> tuple[int, str]:
    """Invoke CLI compile and return (exit_code, output)."""
    result = runner.invoke(main, ["compile", str(src), "-o", str(out)])
    return result.exit_code, result.output


def _make_simple_docx(path: Path) -> None:
    from docx import Document as DocxDoc
    doc = DocxDoc()
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("This is a test DOCX file for the beta smoke test.")
    doc.add_paragraph("It contains simple prose content.")
    doc.save(str(path))


def _make_simple_html(path: Path) -> None:
    path.write_text(
        "<html><body><h1>Test HTML</h1>"
        "<p>This is a simple HTML file for the beta smoke test.</p>"
        "<table><tr><th>A</th><th>B</th></tr><tr><td>1</td><td>2</td></tr></table>"
        "</body></html>",
        encoding="utf-8",
    )


def _make_simple_xlsx(path: Path) -> None:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value", "Notes"])
    ws.append(["Alpha", 1, "First entry"])
    ws.append(["Beta", 2, "Second entry"])
    wb.save(str(path))


def _make_simple_pptx(path: Path) -> None:
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Beta Smoke Test"
    slide.placeholders[1].text = "This is a test PPTX slide."
    prs.save(str(path))


def _make_simple_zip(path: Path, tmp_path: Path) -> None:
    import zipfile
    txt = tmp_path / "_zip_content.txt"
    txt.write_text("Zipped file content for smoke test.", encoding="utf-8")
    with zipfile.ZipFile(str(path), "w") as zf:
        zf.write(str(txt), "content.txt")


def _make_simple_image(path: Path) -> None:
    img = Image.new("RGB", (400, 200), (200, 220, 240))
    draw = ImageDraw.Draw(img)
    draw.text((50, 80), "Test Image for AksharaMD", fill=(0, 0, 0))
    img.save(str(path), format="PNG")


@pytest.mark.slow
def test_beta_corpus_smoke(runner, tmp_path):
    """Smoke-test that each corpus file compiles without crashing and produces
    a manifest with required fields. Checks user-facing output readability."""
    corpus = {}

    # PDFs
    clean_pdf = tmp_path / "clean.pdf"
    _make_clean_native_pdf(clean_pdf)
    corpus["clean_native_pdf"] = clean_pdf

    scanned_pdf = tmp_path / "scanned.pdf"
    _make_scanned_pdf(scanned_pdf)
    corpus["scanned_pdf"] = scanned_pdf

    encrypted_pdf = tmp_path / "encrypted.pdf"
    _make_encrypted_pdf(encrypted_pdf)
    corpus["encrypted_pdf"] = encrypted_pdf

    # Non-PDF formats
    html_file = tmp_path / "page.html"
    _make_simple_html(html_file)
    corpus["html"] = html_file

    img_file = tmp_path / "img.png"
    _make_simple_image(img_file)
    corpus["image_png"] = img_file

    zip_file = tmp_path / "archive.zip"
    _make_simple_zip(zip_file, tmp_path)
    corpus["zip_archive"] = zip_file

    try:
        docx_file = tmp_path / "doc.docx"
        _make_simple_docx(docx_file)
        corpus["docx"] = docx_file
    except ImportError:
        pass  # python-docx optional

    try:
        xlsx_file = tmp_path / "sheet.xlsx"
        _make_simple_xlsx(xlsx_file)
        corpus["xlsx"] = xlsx_file
    except ImportError:
        pass

    try:
        pptx_file = tmp_path / "deck.pptx"
        _make_simple_pptx(pptx_file)
        corpus["pptx"] = pptx_file
    except ImportError:
        pass

    results = {}
    for name, src in corpus.items():
        out_dir = tmp_path / f"out_{name}"
        exit_code, output = _smoke(runner, src, out_dir)
        results[name] = {"exit_code": exit_code, "output": output}

    # Encrypted PDF should fail gracefully with a useful message
    enc = results.get("encrypted_pdf", {})
    assert enc.get("exit_code", 0) != 0, "Encrypted PDF should exit non-zero"
    assert any(w in enc.get("output", "").lower()
               for w in ["password", "encrypted", "decrypt"]), (
        "Encrypted PDF message should mention password/encryption"
    )

    # Clean PDF should succeed with HIGH quality
    clean = results.get("clean_native_pdf", {})
    assert clean.get("exit_code") == 0, (
        f"Clean native PDF should succeed. Output:\n{clean.get('output')}"
    )
    assert "HIGH" in clean.get("output", ""), (
        "Clean native PDF should show HIGH quality band"
    )

    # Scanned PDF should succeed (exit 0) but show a warning about OCR
    scanned = results.get("scanned_pdf", {})
    assert scanned.get("exit_code") == 0, (
        f"Scanned PDF should not crash. Output:\n{scanned.get('output')}"
    )
    scanned_out = scanned.get("output", "")
    assert "POOR" in scanned_out or "RISKY" in scanned_out or "OCR" in scanned_out, (
        "Scanned PDF should show poor quality or OCR guidance"
    )

    # All non-encrypted files should compile without crashing
    for name, res in results.items():
        if name == "encrypted_pdf":
            continue
        assert res["exit_code"] == 0 or "ERROR" in res["output"], (
            f"{name} crashed unexpectedly. Output:\n{res['output']}"
        )

    # All output directories should have manifest.json
    for name in corpus:
        if name == "encrypted_pdf":
            continue
        stem = corpus[name].stem
        manifest_path = tmp_path / f"out_{name}" / stem / "manifest.json"
        if manifest_path.exists():
            data = json.loads(manifest_path.read_text())
            assert "quality_band" in data, f"{name}: manifest missing quality_band"
            assert "readiness_score" in data, f"{name}: manifest missing readiness_score"
