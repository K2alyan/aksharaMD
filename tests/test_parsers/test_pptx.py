from __future__ import annotations

import io
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx", reason="python-pptx not installed")

from pptx import Presentation

from aksharamd.compiler import Compiler
from aksharamd.models.block import BlockType


def _make_pptx(tmp_path: Path, slides: list[dict]) -> Path:
    """Create a PPTX file from a list of slide specs."""
    prs = Presentation()
    for spec in slides:
        layout_idx = spec.get("layout", 1)
        try:
            layout = prs.slide_layouts[layout_idx]
        except IndexError:
            layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        if "title" in spec:
            try:
                slide.shapes.title.text = spec["title"]
            except AttributeError:
                pass

        if "body" in spec:
            try:
                slide.placeholders[1].text = spec["body"]
            except (IndexError, AttributeError, KeyError):
                pass

    path = tmp_path / "presentation.pptx"
    buf = io.BytesIO()
    prs.save(buf)
    path.write_bytes(buf.getvalue())
    return path


def _compile(path: Path, tmp_path: Path):
    return Compiler(output_dir=str(tmp_path / "out")).compile(str(path))


# ── basic parsing ─────────────────────────────────────────────────────────────

def test_pptx_produces_document(tmp_path):
    p = _make_pptx(tmp_path, [{"title": "Hello", "body": "World"}])
    ctx = _compile(p, tmp_path)
    assert ctx.document is not None
    assert ctx.document.file_type == "pptx"


def test_pptx_pages_match_slide_count(tmp_path):
    p = _make_pptx(tmp_path, [
        {"title": "Slide One", "body": "Content 1"},
        {"title": "Slide Two", "body": "Content 2"},
        {"title": "Slide Three", "body": "Content 3"},
    ])
    ctx = _compile(p, tmp_path)
    assert ctx.document is not None
    assert ctx.document.pages == 3


def test_pptx_title_extracted(tmp_path):
    p = _make_pptx(tmp_path, [{"title": "My Presentation Title", "body": "Some content."}])
    ctx = _compile(p, tmp_path)
    assert ctx.document is not None
    # Title should come from first slide title
    assert ctx.document.title is not None


def test_pptx_has_blocks(tmp_path):
    p = _make_pptx(tmp_path, [{"title": "Intro", "body": "Introduction paragraph text."}])
    ctx = _compile(p, tmp_path)
    assert len(ctx.document.blocks) > 0


def test_pptx_document_metadata_contains_slide_count(tmp_path):
    p = _make_pptx(tmp_path, [{"title": "Test"}])
    ctx = _compile(p, tmp_path)
    assert ctx.document.metadata["slides"] >= 1


def test_pptx_slide_heading_block_present(tmp_path):
    p = _make_pptx(tmp_path, [{"title": "Section Title", "body": "Body text"}])
    ctx = _compile(p, tmp_path)
    headings = [b for b in ctx.document.blocks if b.type == BlockType.HEADING]
    assert len(headings) >= 1


def test_pptx_empty_presentation(tmp_path):
    p = _make_pptx(tmp_path, [])
    ctx = _compile(p, tmp_path)
    assert ctx.document is not None


def test_pptx_multiple_slides_content(tmp_path):
    p = _make_pptx(tmp_path, [
        {"title": "Agenda", "body": "Overview of topics."},
        {"title": "Background", "body": "History and context."},
    ])
    ctx = _compile(p, tmp_path)
    all_content = " ".join(b.content for b in ctx.document.blocks)
    assert "Agenda" in all_content or "Background" in all_content


def test_pptx_corrupt_file_does_not_crash(tmp_path):
    f = tmp_path / "corrupt.pptx"
    f.write_bytes(b"PK\x03\x04" + b"\x00" * 50)
    ctx = _compile(f, tmp_path)
    assert ctx is not None
